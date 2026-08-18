"""Generate the Graxella capabilities deck — "Capabilities in Action".

Run:  python build_capabilities_ppt.py
Output: Graxella_Capabilities_Deck.pptx (in this folder).

The companion to the flagship deck: what the graxella PACKAGE can do
today, shown through scenarios and user stories, with an explicit
"as-is vs graxella-equipped" comparison on every capability and a
master comparison table. Every number in here is a measured number
from the repo's test suites and the real-LLM showcase (08_refund_desk,
qwen2.5:3b).

Uses python-pptx (>=1.0). All colours + fonts set explicitly so the deck
renders identically regardless of the machine's PowerPoint theme.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# --- palette (evidence-ledger identity: paper / ink / verdigris / amber) ---
PAPER = RGBColor(0xFA, 0xFB, 0xF9)
PANEL = RGBColor(0xF0, 0xF4, 0xF0)
INK = RGBColor(0x1C, 0x25, 0x21)
MUTED = RGBColor(0x5C, 0x6A, 0x63)
GREEN = RGBColor(0x14, 0x80, 0x5E)   # verdigris accent — the ledger
AMBER = RGBColor(0xA8, 0x69, 0x1C)   # enterprise pain markers
RED = RGBColor(0xA9, 0x43, 0x2E)     # failure / as-is column
LINE = RGBColor(0xDC, 0xE3, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_TINT = RGBColor(0xE2, 0xEF, 0xE9)
RED_TINT = RGBColor(0xF7, 0xEA, 0xE6)
DARK_GREEN_ON_INK = RGBColor(0x43, 0xC4, 0x95)

HEAD_FONT = "Segoe UI"
BODY_FONT = "Segoe UI"
MONO_FONT = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ------------------------------------------------------------------ helpers
def new_slide():
    return prs.slides.add_slide(BLANK)


def add_bg(slide, colour=PAPER):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = colour
    bg.shadow.inherit = False
    return bg


def _set(p, text, *, size, colour, font=BODY_FONT, bold=False,
         italic=False, align=None, spacing=None):
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = colour
    p.font.name = font
    if align is not None:
        p.alignment = align
    if spacing is not None:
        p.space_after = Pt(spacing)
    return p


def add_eyebrow(slide, text, *, y=0.42, colour=GREEN, x=0.7, w=12.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.32))
    _set(tb.text_frame.paragraphs[0], text.upper(), size=12, colour=colour,
         font=MONO_FONT, bold=True)
    return tb


def add_title(slide, text, *, y=0.72, size=30, colour=INK, x=0.7, w=12.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.85))
    _set(tb.text_frame.paragraphs[0], text, size=size, colour=colour,
         font=HEAD_FONT, bold=True)
    return tb


def add_sub(slide, text, *, y=1.42, size=14, colour=MUTED, x=0.7, w=12.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    _set(tf.paragraphs[0], text, size=size, colour=colour)
    return tb


def add_bullets(slide, items, *, x=0.7, y=2.05, w=11.9, h=4.9,
                size=14, spacing=10, head_colour=GREEN, body_colour=INK):
    """items: list of str or (head, body) tuples."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        if isinstance(item, tuple):
            head, body = item
            run = p.add_run()
            run.text = f"{head} — "
            run.font.bold = True
            run.font.size = Pt(size)
            run.font.color.rgb = head_colour
            run.font.name = HEAD_FONT
            run2 = p.add_run()
            run2.text = body
            run2.font.size = Pt(size)
            run2.font.color.rgb = body_colour
            run2.font.name = BODY_FONT
        else:
            _set(p, item, size=size, colour=body_colour)
    return tb


def add_panel(slide, x, y, w, h, *, fill=PANEL, line=LINE, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    if radius:
        try:
            sh.adjustments[0] = 0.06
        except Exception:
            pass
    return sh


def panel_text(sh, lines, *, size=12, colour=INK, font=MONO_FONT,
               align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=4,
               margin=0.14):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin * 0.7)
    tf.margin_bottom = Inches(margin * 0.7)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, tuple):
            text, colr, bold = line
            _set(p, text, size=size, colour=colr, font=font, bold=bold,
                 align=align, spacing=spacing)
        else:
            _set(p, line, size=size, colour=colour, font=font,
                 align=align, spacing=spacing)
    return sh


def add_footer(slide, text="GRAXELLA · CAPABILITIES IN ACTION"):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(7.08),
                                  Inches(12.0), Inches(0.3))
    _set(tb.text_frame.paragraphs[0], text, size=9, colour=MUTED,
         font=MONO_FONT)
    return tb


def add_table(slide, rows, *, x=0.7, y=2.0, w=11.9, col_widths=None,
              row_h=0.52, header_h=0.42, size=12, header_colour=GREEN):
    n_rows, n_cols = len(rows), len(rows[0])
    total_h = header_h + row_h * (n_rows - 1)
    gfx = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                 Inches(w), Inches(total_h))
    table = gfx.table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            table.columns[ci].width = Inches(cw)
    table.rows[0].height = Inches(header_h)
    for ri in range(1, n_rows):
        table.rows[ri].height = Inches(row_h)
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (GREEN_TINT if ri == 0 else
                                        (PAPER if ri % 2 else PANEL))
            cell.margin_left = Inches(0.09)
            cell.margin_right = Inches(0.09)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if ri == 0:
                _set(p, cell_text.upper(), size=size - 1, colour=header_colour,
                     font=MONO_FONT, bold=True)
            else:
                _set(p, cell_text, size=size, colour=INK, bold=(ci == 0))
    return gfx


def add_compare(slide, asis_lines, equipped_lines, *, y=4.55, h=2.15,
                size=11.5):
    """The signature device of this deck: AS-IS (red) vs EQUIPPED (green)."""
    left = add_panel(slide, 0.7, y, 5.85, h, fill=RED_TINT,
                     line=RGBColor(0xE3, 0xC6, 0xBC), radius=True)
    panel_text(left, [("AS-IS (FRAMEWORK ALONE)", RED, True)] + asis_lines,
               size=size, font=BODY_FONT, spacing=3)
    right = add_panel(slide, 6.75, y, 5.85, h, fill=GREEN_TINT,
                      line=RGBColor(0xBE, 0xD9, 0xCC), radius=True)
    panel_text(right, [("GRAXELLA-EQUIPPED", GREEN, True)] + equipped_lines,
               size=size, font=BODY_FONT, spacing=3)
    return left, right


# =========================================================== 1 · TITLE
s = new_slide()
add_bg(s, INK)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.4))
_set(tb.text_frame.paragraphs[0], "GRAXELLA · THE PACKAGE", size=14,
     colour=DARK_GREEN_ON_INK, font=MONO_FONT, bold=True)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.3))
_set(tb.text_frame.paragraphs[0], "Capabilities in Action", size=54,
     colour=WHITE, font=HEAD_FONT, bold=True)
tb = s.shapes.add_textbox(Inches(0.9), Inches(3.8), Inches(11.0), Inches(0.9))
tf = tb.text_frame
tf.word_wrap = True
_set(tf.paragraphs[0],
     "What the graxella package does today — seven capabilities, three "
     "scenarios, six user stories, and an honest as-is vs graxella-equipped "
     "comparison on every one of them.",
     size=19, colour=RGBColor(0xC9, 0xD4, 0xCE))
tb = s.shapes.add_textbox(Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.4))
_set(tb.text_frame.paragraphs[0],
     "v0.1 · 2026-08 · EVERY NUMBER MEASURED · 354 TESTS · REAL-LLM SHOWCASE"
     " ON QWEN2.5:3B",
     size=10, colour=RGBColor(0x7E, 0x8D, 0x85), font=MONO_FONT)

# =================================== 2 · WHAT IT IS + THE ONE-LINE ADOPTION
s = new_slide()
add_bg(s)
add_eyebrow(s, "What graxella is")
add_title(s, "A learning substrate UNDER your agent framework")
add_sub(s, "Keep LangGraph / LangChain and your agents exactly as they are. "
           "One line puts the evidence loop underneath them.")
code = add_panel(s, 0.7, 2.0, 7.3, 3.1, fill=INK, line=INK, radius=True)
panel_text(code, [
    ('llm = ChatOllama(model="qwen2.5:3b")   # any model', MUTED, False),
    ("triage    = create_react_agent(llm, [check_order,", WHITE, False),
    ("                lookup_policy, shipping], name=\"triage\")", WHITE, False),
    ("responder = create_react_agent(llm, [send_email],", WHITE, False),
    ("                               name=\"responder\")", WHITE, False),
    ("", WHITE, False),
    ("app = graxella.mesh([triage, responder],", DARK_GREEN_ON_INK, True),
    ("                    domain=\"refunds\")   # the new line",
     DARK_GREEN_ON_INK, True),
    ("t = app.run_trajectory(\"order 1042 damaged - refund?\")",
     WHITE, False),
], size=13)
side = add_panel(s, 8.25, 2.0, 4.35, 3.1, radius=True)
panel_text(side, [
    ("WHAT JUST TURNED ON", GREEN, True),
    ("· every dispatch = decision + typed outcome", INK, False),
    ("· case recall from verified history", INK, False),
    ("· bounded multi-hop with loop detection", INK, False),
    ("· tool heal ladder + evidence gate", INK, False),
    ("· cited audit trail, topology, OTel", INK, False),
], size=12, font=BODY_FONT, spacing=6)
add_bullets(s, [
    ("Model-agnostic by construction", "graxella never sees the model choice; "
     "all evidence is scoped per model_id, so learnings never leak across "
     "models."),
    ("Framework-agnostic by construction", "agents are plain callables or "
     "LangGraph graphs; tools adapt in one lambda. No YAML, no handoff "
     "protocol, no schema to author."),
    ("Silent plumbing, loud failures", "nothing about your flow changes — "
     "but every dead end becomes a signal a human can see, never a silent "
     "retry."),
], y=5.35, size=13, spacing=6)
add_footer(s)

# ============================================== 3 · THE CAPABILITY MAP
s = new_slide()
add_bg(s)
add_eyebrow(s, "The capability map")
add_title(s, "Seven capabilities, one loop: evidence in, reliability out")
add_sub(s, "Each one ships in the package today, with tests. The next seven "
           "slides take them one at a time.")
caps = [
    ("1 · OUTCOME LEDGER", "Every dispatch auto-records a decision and a "
     "typed outcome — mnema-backed, queryable, immutable."),
    ("2 · CASE RECALL", "Verified past outcomes injected into context — "
     "agents stop repeating solved mistakes."),
    ("3 · TOOL SELF-HEALING", "The heal ladder: fail once on drift, learn "
     "forever. LLM appears exactly once."),
    ("4 · EVIDENCE GATE", "No behavior change without cited evidence — "
     "Bayesian, self-calibrating, human-reviewed when cold."),
    ("5 · BOUNDED TRAJECTORIES", "Multi-hop with loop detection, budgets, "
     "typed handoffs. Runaways escalate, never burn."),
    ("6 · PROGRESSIVE DISCLOSURE", "Agent capabilities revealed L0 to L3 as "
     "needed — token spend follows relevance."),
    ("7 · TRUST + OBSERVABILITY", "Cited tool trust, gate.why(), OTel "
     "spans, live topology map, control-plane API + CLI."),
]
for i, (head, body) in enumerate(caps):
    col, row = i % 4, i // 4
    px, py = 0.7 + col * 3.075, 2.05 + row * 2.3
    p = add_panel(s, px, py, 2.87, 2.1, radius=True)
    panel_text(p, [(head, GREEN, True), (body, INK, False)],
               size=11, font=BODY_FONT, spacing=5)
p = add_panel(s, 9.925, 4.35, 2.87, 2.1, fill=INK, line=INK, radius=True)
panel_text(p, [("THE LOOP", DARK_GREEN_ON_INK, True),
               ("dispatch → outcome → ledger → proposal → gate → promotion. "
                "Volume becomes reliability.", WHITE, False)],
           size=11, font=BODY_FONT, spacing=5)
add_footer(s)

# ====================================== 4 · CAPABILITY 1 — OUTCOME LEDGER
s = new_slide()
add_bg(s)
add_eyebrow(s, "Capability 1 of 7 · outcome ledger")
add_title(s, "Every dispatch leaves evidence — with zero developer effort")
add_sub(s, "route() records the decision (task, chosen agent, domain, "
           "model_id) and the typed outcome (ok, error class, tokens, "
           "latency, tools used) into the mnema ledger. Hot path <2ms via a "
           "write-ahead buffer.")
add_bullets(s, [
    ("Typed, not textual", "outcomes are structured records, not log lines — "
     "outcome_stats() answers ok-rate, token and latency questions in one "
     "call."),
    ("Immutable + cited", "ledger assertions are append-only with "
     "provenance; trajectories cite their hop decisions (derived_from)."),
    ("Scoped per model", "every record carries model_id — swap models and "
     "the evidence never cross-contaminates."),
], y=2.15, size=13, spacing=7)
add_compare(s, [
    ("Telemetry is unstructured logs; \"how often does triage fail on "
     "refunds?\" means an hour of grep.", INK, False),
    ("Sessions vanish when the process exits; nobody can reconstruct why "
     "an agent was chosen.", INK, False),
], [
    ("memory.outcome_stats() — count, ok_rate, tokens, latency — one call, "
     "any slice.", INK, False),
    ("Every choice is a decision record with an explanation; every result "
     "is a typed outcome citing it.", INK, False),
])
add_footer(s)

# ======================================== 5 · CAPABILITY 2 — CASE RECALL
s = new_slide()
add_bg(s)
add_eyebrow(s, "Capability 2 of 7 · case recall")
add_title(s, "Agents that remember what actually worked")
add_sub(s, "Before dispatch, graxella retrieves similar past cases whose "
           "outcomes are VERIFIED in the ledger and injects them as a "
           "compact context block (the Memento pattern).")
add_bullets(s, [
    ("Grounded in outcomes, not vibes", "only cases with recorded verified "
     "outcomes are recalled — no hallucinated \"experience\"."),
    ("Automatic", "recall is on by default in mesh(); the agent prompt "
     "gains a short evidence block, nothing else changes."),
    ("Proven live", "in the refund-desk showcase, request 3 saw requests "
     "1–2 as verified experience (recall.injected tracer event)."),
], y=2.15, size=13, spacing=7)
add_compare(s, [
    ("Every session starts from zero; the same mistake is re-made and "
     "re-billed daily.", INK, False),
    ("\"Memory\" add-ons store chat text — unverified, unscoped, and "
     "unauditable.", INK, False),
], [
    ("Solved cases compound: month six is measurably more reliable than "
     "week one.", INK, False),
    ("Recall blocks cite ledger assertions — you can audit exactly what "
     "the agent was reminded of.", INK, False),
])
add_footer(s)

# ==================================== 6 · CAPABILITY 3 — TOOL SELF-HEALING
s = new_slide()
add_bg(s)
add_eyebrow(s, "Capability 3 of 7 · tool self-healing")
add_title(s, "The heal ladder: fail once, learn forever")
add_sub(s, "ToolInterceptor wraps any tool (MCP, HTTP, LangChain — one "
           "lambda). On schema drift it climbs the cheapest rung that "
           "works. The LLM appears at exactly one rung, exactly once.")
rungs = [
    ("RUNG 1 · HAPPY PATH", "primary succeeds; typed outcome recorded.",
     PANEL, GREEN),
    ("RUNG 2 · PROMOTED HEAL", "a gate-approved transform from the "
     "Rulebook. ZERO LLM. Warms the gate for the next proposal.",
     GREEN_TINT, GREEN),
    ("RUNG 2.5 · PROPOSED RECIPE", "recipe awaiting review is reused "
     "deterministically — the LLM truly fires once per process.",
     GREEN_TINT, GREEN),
    ("RUNG 3 · HEAL ONCE", "the pluggable healer (an LLM, once) proposes a "
     "TransformRecipe; success ships as a gated Proposal with paired-replay "
     "evidence.", PANEL, AMBER),
    ("RUNG 4 · LOUD FAILURE", "nothing worked: typed failure outcome, "
     "re-raise. Never a silent retry.", RED_TINT, RED),
]
for i, (head, body, fill, hcol) in enumerate(rungs):
    p = add_panel(s, 0.7 + i * 2.44, 2.15, 2.32, 2.15, fill=fill,
                  radius=True)
    panel_text(p, [(head, hcol, True), (body, INK, False)], size=10,
               font=BODY_FONT, spacing=4)
add_compare(s, [
    ("A renamed API field breaks every call until a human ships a code "
     "fix — or an LLM \"fixes\" it on every single call, silently, at "
     "full token price.", INK, False),
], [
    ("Measured in the real-LLM showcase: two drift events, LLM healer "
     "invocations total = 1. The second heal was a promoted, cited, "
     "deterministic transform.", INK, False),
])
add_footer(s)

# ====================================== 7 · CAPABILITY 4 — EVIDENCE GATE
s = new_slide()
add_bg(s)
add_eyebrow(s, "Capability 4 of 7 · evidence gate")
add_title(s, "Nothing changes runtime behavior without cited evidence")
add_sub(s, "Every learned artifact (transform, rule, binding) is a "
           "Proposal decided by a Bayesian gate per (domain, kind, target, "
           "model) tuple — never by a score someone tuned, never by an LLM.")
add_bullets(s, [
    ("Self-calibrating threshold", "thr(n) = 0.85 + 0.10·e^(−n/20): strict "
     "when evidence is thin, earned confidence as it accumulates."),
    ("Cold start goes to a human", "unproven tuples land in the review "
     "queue (NEEDS_HUMAN); approvals are themselves ledger assertions."),
    ("Constitution beats everyone", "hard blocks outrank all evidence — "
     "including human approvals. Verdicts render with citations via "
     "gate.why()."),
    ("Versioned + rollbackable", "promotions carry lineage (gate verdict, "
     "approver); rollback is a state change, not an archaeology dig."),
], y=2.15, size=13, spacing=7)
add_compare(s, [
    ("\"Learning\" = someone edits a prompt on a Friday; no record, no "
     "rollback, no threshold.", INK, False),
    ("Or: autonomous self-modification you can't explain to an auditor.",
     INK, False),
], [
    ("Every behavior change: proposed, evidence-cited, gate-decided, "
     "human-reviewed when cold, promoted with lineage.", INK, False),
    ("gate.why() prints the full verdict chain — audit-ready by "
     "construction.", INK, False),
])
add_footer(s)

# ================================ 8 · CAPABILITY 5 — BOUNDED TRAJECTORIES
s = new_slide()
add_bg(s)
add_eyebrow(s, "Capability 5 of 7 · bounded trajectories")
add_title(s, "Multi-hop that cannot run away")
add_sub(s, "run_trajectory() drives route → execute → assess → handoff "
           "with four containments aimed at MAST's two most frequent "
           "failure modes: step repetition (15.7%) and no termination "
           "(12.4%).")
add_bullets(s, [
    ("Loop detection", "a repeated (agent, response-state) signature stops "
     "the chain and escalates — never re-runs."),
    ("Budgets", "max hops, tokens, wallclock; exhaustion escalates with a "
     "ledger signal instead of burning the bill."),
    ("Typed handoffs", "agents hand off with 'HANDOFF: <peer> :: <task>' — "
     "every hop is a full audited route(). Self-handoffs and unknown peers "
     "complete loudly (a real-qwen2.5:3b lesson, now a test)."),
    ("The chain is a ledger object", "predicate=\"trajectory\", "
     "derived_from = every hop decision — chain-level evidence for miners."),
], y=2.15, size=13, spacing=7)
add_compare(s, [
    ("Two agents politely hand the same task back and forth until the "
     "token budget — or the credit card — gives out.", INK, False),
    ("A 3B model hands off to itself forever; nobody notices until the "
     "invoice.", INK, False),
], [
    ("Loop detected at the first repeat; contained in 2–3 hops; a "
     "trajectory_escalation signal reaches a human.", INK, False),
    ("Self-handoff completes immediately with a loud tracer event — "
     "measured, tested, shipped.", INK, False),
])
add_footer(s)

# ============================ 9 · CAPABILITY 6 — PROGRESSIVE DISCLOSURE
s = new_slide()
add_bg(s)
add_eyebrow(s, "Capability 6 of 7 · progressive disclosure")
add_title(s, "Token spend follows relevance, not roster size")
add_sub(s, "Peer capabilities are revealed in levels — the router decides "
           "who earns more tokens. No more pasting every agent card into "
           "every prompt.")
levels = [
    ("L0 · AWARENESS", "≤ 25 tokens", "every peer, always: name + "
     "one-line skill."),
    ("L1 · SHORTLIST", "≤ 120 tokens", "route-relevant peers get a skill "
     "digest, pre-routed by the deterministic router."),
    ("L2 · WINNER", "≤ 600 tokens", "the chosen peer's full card — "
     "capabilities, constraints, examples."),
    ("L3 · INVOCATION", "as needed", "live handoff context for the actual "
     "call."),
]
for i, (head, budget, body) in enumerate(levels):
    p = add_panel(s, 0.7 + i * 3.075, 2.15, 2.87, 2.0,
                  fill=GREEN_TINT if i < 2 else PANEL, radius=True)
    panel_text(p, [(head, GREEN, True), (budget, AMBER, True),
                   (body, INK, False)], size=11, font=BODY_FONT, spacing=4)
add_compare(s, [
    ("50 agents × full capability cards × every prompt = thousands of "
     "\"waiting tokens\" per call, paid forever.", INK, False),
], [
    ("Deterministic TF-IDF routing picks the shortlist BEFORE the LLM "
     "sees anything; disclosure budgets are enforced, not aspirational.",
     INK, False),
])
add_footer(s)

# ========================== 10 · CAPABILITY 7 — TRUST + OBSERVABILITY
s = new_slide()
add_bg(s)
add_eyebrow(s, "Capability 7 of 7 · trust + observability")
add_title(s, "Every answer to \"why?\" is one call away")
add_sub(s, "The same ledger that powers learning powers explanation — "
           "no second telemetry system to keep honest.")
surfaces = [
    ("gate.why(id)", "the cited verdict chain for any learned behavior — "
     "evidence, threshold, approver, constitution checks."),
    ("tool_trust()", "Laplace-smoothed trust per tool from real outcomes: "
     "(s+1)/(s+f+2) — no hand-tuned scores."),
    ("OTel bridge", "gen_ai.* spans + graxella.* attributes — traces land "
     "in Jaeger/your APM, not a proprietary viewer."),
    ("Topology map", "one HTML file, zero dependencies: agents, tools, "
     "trust and flow at a glance."),
    ("Control plane", "FastAPI service: /health /stats /trust /gate/* "
     "/ingest — plus a CLI (graxella tokens, graxella gate)."),
]
for i, (head, body) in enumerate(surfaces):
    col, row = i % 3, i // 3
    p = add_panel(s, 0.7 + col * 4.1, 2.1 + row * 1.6, 3.9, 1.45,
                  radius=True)
    panel_text(p, [(head, GREEN, True), (body, INK, False)], size=11,
               font=BODY_FONT, spacing=4)
add_compare(s, [
    ("\"Why did the agent do that?\" → grep logs, interview the developer, "
     "shrug.", INK, False),
], [
    ("Decisions, outcomes, trust, verdicts and topology are queryable "
     "surfaces of one immutable ledger.", INK, False),
], y=5.45, h=1.5)
add_footer(s)

# ================================ 11 · SCENARIO A — THE REFUND DESK (REAL)
s = new_slide()
add_bg(s)
add_eyebrow(s, "Scenario A · measured, not simulated", colour=AMBER)
add_title(s, "The refund desk — two real qwen2.5:3b agents, live")
add_sub(s, "showcase/08_refund_desk.py: TRIAGE (order/policy/shipping "
           "tools) hands off to RESPONDER (email) over a drifted internal "
           "shipping API. Everything below is from the actual run.")
p = add_panel(s, 0.7, 2.1, 5.9, 4.3, fill=INK, line=INK, radius=True)
panel_text(p, [
    ("$ uv run python showcase/08_refund_desk.py", MUTED, False),
    ("[1] order 1042 damaged, wants refund...", WHITE, False),
    ("    hops=['triage','responder']  completed (6s)", WHITE, False),
    ("    reply: 'refundable within 30 days...'", DARK_GREEN_ON_INK, False),
    ("[2] refund for order 2077 desk lamp...", WHITE, False),
    ("    completed (4s) — 40 days old: store credit", WHITE, False),
    ("[3] order 1042 damaged again...", WHITE, False),
    ("    case-recall: saw [1]-[2] as experience", DARK_GREEN_ON_INK, False),
    ("", WHITE, False),
    ("shipping drift: HTTP_410_GONE schema deprecated", RED, False),
    ("heal → gate → operator approve → promote", WHITE, False),
    ("LLM healer invocations, total: 1", DARK_GREEN_ON_INK, True),
], size=12)
add_bullets(s, [
    ("Answers in 4–6 s", "with correct policy applied (30-day damaged "
     "window vs store credit) — small model, contained runtime."),
    ("One LLM heal, ever", "drift #1 pays the LLM once; drift #2 heals via "
     "the promoted transform at zero LLM cost."),
    ("Full audit for free", "every hop a decision+outcome; gate.why() "
     "cites the operator's approval; topology.html written."),
    ("Model honesty", "qwen2.5:3b sometimes mangles handoff markers — the "
     "runtime contains it loudly instead of spinning."),
], x=6.9, y=2.1, w=5.7, size=12.5, spacing=7)
add_footer(s)

# ================================== 12 · SCENARIO B — API DRIFT AT 2 AM
s = new_slide()
add_bg(s)
add_eyebrow(s, "Scenario B · user story: the on-call engineer", colour=AMBER)
add_title(s, "2:00 AM: a vendor renames a field")
add_sub(s, "\"As an on-call SRE, I want tool schema drift to heal itself "
           "under governance, so that drift stops paging me.\"")
steps = [
    ("02:00", "vendor ships shipping.v2; order_id becomes order_ref. "
     "Every agent call starts failing with HTTP_410_GONE."),
    ("02:00:04", "the interceptor recognizes the drift signature, invokes "
     "the healer LLM ONCE; the proposed rename works on replay."),
    ("02:00:05", "a Proposal with paired-replay evidence lands in the gate. "
     "Cold tuple → review queue. Calls keep flowing on the proposed "
     "recipe (rung 2.5, zero further LLM)."),
    ("09:00", "a human reviews the one-line diff over coffee, approves; "
     "the transform is promoted with lineage. The tuple is warmer for "
     "next time."),
]
for i, (t, body) in enumerate(steps):
    p = add_panel(s, 0.7, 2.15 + i * 0.98, 11.9, 0.88,
                  fill=PANEL if i % 2 else PAPER, radius=True)
    panel_text(p, [(f"{t} — {body}", INK, False)], size=12.5,
               font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)
add_compare(s, [
    ("Pager fires at 02:00. Agents fail all night, or a runtime LLM "
     "\"fixes\" every call silently at full price with no record.",
     INK, False),
], [
    ("Zero pages. One LLM call, one human review, one promoted transform "
     "— cited, versioned, rollbackable.", INK, False),
], y=6.15, h=1.05)
add_footer(s)

# ============================== 13 · SCENARIO C — THE RUNAWAY CHAIN
s = new_slide()
add_bg(s)
add_eyebrow(s, "Scenario C · user story: the ops lead", colour=AMBER)
add_title(s, "Two polite agents, one infinite bill")
add_sub(s, "\"As a support-ops lead, I want runaway agent chains contained "
           "automatically, so that a loop never becomes an invoice.\"")
add_bullets(s, [
    ("The setup", "PING hands the task to PONG; PONG hands it straight "
     "back. Both are being perfectly \"helpful\". This is MAST FM-1.3 — "
     "the single most frequent multi-agent failure mode observed (15.7%)."),
    ("What graxella does", "hop 1: ping. Hop 2: pong. Hop 3: the (agent, "
     "response-state) signature repeats → status=loop_detected, chain "
     "stopped, trajectory_escalation signal recorded, tracer event "
     "emitted. A human sees it; the bill doesn't."),
    ("The subtler cousin", "a small model handing off to itself — seen "
     "live with qwen2.5:3b. The runtime completes the work loudly "
     "(self_handoff_ignored event) instead of spinning. Found in the real "
     "showcase, fixed in the substrate, locked with a test."),
    ("Budgets as the backstop", "TrajectoryBudget(max_hops, max_tokens, "
     "max_wallclock_s) — even novel loop shapes exhaust a budget and "
     "escalate instead of running away."),
], y=2.1, size=13, spacing=8)
add_compare(s, [
    ("Discovered days later as a line item on the token invoice.",
     INK, False),
], [
    ("Contained within 2–3 hops; escalated as a typed signal with the "
     "full hop history in the ledger.", INK, False),
], y=5.75, h=1.2)
add_footer(s)

# ======================================= 14 · USER STORIES — BUILD SIDE
s = new_slide()
add_bg(s)
add_eyebrow(s, "User stories · the people who build")
add_title(s, "Stories the package already delivers — build side")
stories = [
    ("PLATFORM ENGINEER",
     "\"As a platform engineer, I want to put a reliability substrate "
     "under our existing LangGraph agents without rewriting them.\"",
     "Shipped as: graxella.mesh([your_agents]) — one line; agents and "
     "tools keep their native shape; no YAML, no protocol to author."),
    ("AGENT DEVELOPER",
     "\"As an agent developer, I want my agents to stop repeating "
     "mistakes they already made last week.\"",
     "Shipped as: automatic case recall — verified outcomes from the "
     "ledger injected as compact context, cited and auditable."),
    ("DATA / ML ENGINEER",
     "\"As an ML engineer, I want to swap or A/B models without losing "
     "or polluting what the system has learned.\"",
     "Shipped as: model_id scoping on every decision, outcome and "
     "promotion — evidence never leaks across models."),
]
for i, (role, want, shipped) in enumerate(stories):
    p = add_panel(s, 0.7, 1.75 + i * 1.75, 11.9, 1.6, radius=True)
    panel_text(p, [(role, GREEN, True), (want, INK, False),
                   (shipped, MUTED, False)], size=12, font=BODY_FONT,
               spacing=4)
add_footer(s)

# ========================================= 15 · USER STORIES — RUN SIDE
s = new_slide()
add_bg(s)
add_eyebrow(s, "User stories · the people who run and answer for it")
add_title(s, "Stories the package already delivers — run side")
stories = [
    ("ON-CALL SRE",
     "\"As an on-call SRE, I want tool drift and agent loops to be "
     "contained and queued for review, so that they stop paging me.\"",
     "Shipped as: the heal ladder (LLM once, then deterministic) + loop "
     "detection + budgets + escalation signals."),
    ("COMPLIANCE OFFICER",
     "\"As a compliance officer, I need every autonomous behavior change "
     "traceable to evidence and an accountable approver.\"",
     "Shipped as: the evidence gate — proposals, cited verdicts, human "
     "approvals as ledger assertions, gate.why() on demand."),
    ("ENGINEERING LEADER",
     "\"As an engineering leader, I want proof the agent program gets "
     "MORE reliable with volume — or I cut it at renewal.\"",
     "Shipped as: outcome_stats() and trust surfaces over the immutable "
     "ledger; the topology map for the exec walkthrough. (The full ROI "
     "value ledger is Phase 4 — next.)"),
]
for i, (role, want, shipped) in enumerate(stories):
    p = add_panel(s, 0.7, 1.75 + i * 1.75, 11.9, 1.6, radius=True)
    panel_text(p, [(role, GREEN, True), (want, INK, False),
                   (shipped, MUTED, False)], size=12, font=BODY_FONT,
               spacing=4)
add_footer(s)

# ================================= 16 · MASTER TABLE — AS-IS VS EQUIPPED
s = new_slide()
add_bg(s)
add_eyebrow(s, "The comparison, in one table")
add_title(s, "As-is vs graxella-equipped", size=28)
add_table(s, [
    ("Situation", "As-is (framework alone)", "Graxella-equipped"),
    ("Tool schema drifts",
     "Fails until a human ships a fix, or pays an LLM on every call",
     "Heal once → gated transform; zero-LLM ever after"),
    ("Agent repeats a solved mistake",
     "Every session starts from zero",
     "Verified cases recalled into context, cited"),
    ("Two agents loop / chain runs away",
     "Token burn until timeout or invoice",
     "Loop detected at first repeat; budgets escalate loudly"),
    ("Behavior needs to change",
     "Prompt edit on a Friday; no trail, no rollback",
     "Proposal → evidence → gate → human when cold → versioned promotion"),
    ("\"Why did it do that?\"",
     "Grep logs, interview the team",
     "gate.why(), decision+outcome ledger, topology map"),
    ("50 agents in one mesh",
     "Every card in every prompt — waiting tokens forever",
     "Progressive disclosure L0→L3; deterministic pre-routing"),
    ("Model gets swapped",
     "Learnings lost or silently wrong",
     "Evidence scoped per model_id; substrate unchanged"),
], y=1.75, col_widths=[3.1, 4.4, 4.4], row_h=0.62, size=11.5)
add_footer(s)

# ========================================== 17 · PROOF TODAY + CLOSE
s = new_slide()
add_bg(s, INK)
add_eyebrow(s, "Where the package stands · 2026-08", colour=DARK_GREEN_ON_INK)
tb = s.shapes.add_textbox(Inches(0.7), Inches(0.85), Inches(12.0),
                          Inches(0.9))
_set(tb.text_frame.paragraphs[0], "Measured, tested, honest", size=32,
     colour=WHITE, font=HEAD_FONT, bold=True)
stats = [
    ("354", "tests green across 4 packages (graxella 125 · mnema 105 · "
     "agent2society 112 · axon-fabric 12)"),
    ("33.5 ms", "route p50 at 1,000 agents (down from 1,174 ms); mesh "
     "build 0.1 s at 1,000 (down from 51 s)"),
    ("8 / 14", "MAST failure modes covered (+3 partial, +2 designed) — "
     "tracked openly, replay-tested, never claimed as live prevention"),
    ("1", "LLM healer invocation, total, across two real drift events in "
     "the qwen2.5:3b showcase"),
]
for i, (n, body) in enumerate(stats):
    col, row = i % 2, i // 2
    p = add_panel(s, 0.7 + col * 6.1, 1.85 + row * 1.75, 5.85, 1.55,
                  fill=RGBColor(0x26, 0x32, 0x2B), line=RGBColor(0x3A, 0x49,
                                                                 0x40),
                  radius=True)
    panel_text(p, [(n, DARK_GREEN_ON_INK, True), (body, WHITE, False)],
               size=13, font=BODY_FONT, spacing=4)
tb = s.shapes.add_textbox(Inches(0.7), Inches(5.55), Inches(12.0),
                          Inches(1.0))
tf = tb.text_frame
tf.word_wrap = True
_set(tf.paragraphs[0], "Fail once. Learn forever.", size=30,
     colour=DARK_GREEN_ON_INK, font=HEAD_FONT, bold=True)
p2 = tf.add_paragraph()
_set(p2, "Next: Phase 4, the Value Ledger — ROI queries, the "
         "reliability-slope dashboard, and the compliance pack that turns "
         "this evidence into the pilot report a VP reads.",
     size=14, colour=RGBColor(0xC9, 0xD4, 0xCE))
tb = s.shapes.add_textbox(Inches(0.7), Inches(7.08), Inches(12.0),
                          Inches(0.3))
_set(tb.text_frame.paragraphs[0],
     "GRAXELLA · CAPABILITIES IN ACTION · EVERY CLAIM TRACES TO A TEST OR "
     "A MEASURED RUN", size=9, colour=RGBColor(0x7E, 0x8D, 0x85),
     font=MONO_FONT)

out = Path(__file__).with_name("Graxella_Capabilities_Deck.pptx")
prs.save(out)
print(f"written: {out.name} ({len(prs.slides._sldIdLst)} slides)")
