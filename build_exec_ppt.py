"""Generate the Graxella executive summary PPT.

Run:  python build_exec_ppt.py
Output: Graxella_Executive_Summary.pptx (in this folder).

Uses python-pptx (>=1.0). All colours + fonts are set explicitly so the
deck looks consistent regardless of the machine's PowerPoint theme.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# --- palette (dark navy + amber accent, easy to read on projectors) -------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
AMBER = RGBColor(0xF6, 0xA5, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE8, 0xEC, 0xF2)
GREY = RGBColor(0x55, 0x60, 0x70)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xC0, 0x39, 0x2B)


# ------------------------------------------------------------------ helpers
def add_bg(slide, colour):
    left, top, width, height = 0, 0, prs.slide_width, prs.slide_height
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = colour
    return bg


def add_title(slide, text, *, colour=WHITE, size=32, y=0.35):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12.4), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = colour
    p.font.name = "Calibri"
    return tb


def add_subtitle(slide, text, *, colour=AMBER, size=16, y=1.15):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12.4), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.italic = True
    p.font.color.rgb = colour
    p.font.name = "Calibri"
    return tb


def add_bullets(slide, items, *, x=0.6, y=1.8, w=12.4, h=5.2,
                size=16, colour=WHITE, spacing=8):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            head, body = item
            p.text = ""
            run = p.add_run()
            run.text = f"{head}  "
            run.font.bold = True
            run.font.size = Pt(size)
            run.font.color.rgb = AMBER
            run.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = body
            run2.font.size = Pt(size)
            run2.font.color.rgb = colour
            run2.font.name = "Calibri"
        else:
            p.text = f"•  {item}"
            p.font.size = Pt(size)
            p.font.color.rgb = colour
            p.font.name = "Calibri"
        p.space_after = Pt(spacing)
    return tb


def add_box(slide, x, y, w, h, *, fill=WHITE, line=NAVY, text="",
            text_colour=NAVY, size=14, bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = text_colour
    p.font.name = "Calibri"
    return shape


def add_arrow(slide, x1, y1, x2, y2, *, colour=AMBER):
    conn = slide.shapes.add_connector(2, Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))  # 2 = STRAIGHT
    conn.line.color.rgb = colour
    conn.line.width = Pt(2.5)
    try:
        # end arrow (best-effort — python-pptx variants differ)
        from pptx.oxml.ns import qn
        line = conn.line._get_or_add_ln()
        tail = line.makeelement(qn("a:tailEnd"), {"type": "triangle"})
        line.append(tail)
    except Exception:
        pass
    return conn


def add_footer(slide, page, total):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.05),
                                  Inches(12.5), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = f"Graxella  ·  Graph-Native Agent Intelligence  ·  {page}/{total}"
    p.font.size = Pt(9)
    p.font.color.rgb = GREY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.RIGHT


# ------------------------------------------------------------------ deck
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


TOTAL = 12
page = [0]


def new(navy=True):
    page[0] += 1
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, NAVY if navy else WHITE)
    add_footer(slide, page[0], TOTAL)
    return slide


# ---- 1. Title -----------------------------------------------------------
s = new()
tb = s.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(12.4), Inches(1.4))
p = tb.text_frame.paragraphs[0]
p.text = "Graxella"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Calibri"

tb = s.shapes.add_textbox(Inches(0.6), Inches(3.7), Inches(12.4), Inches(0.8))
p = tb.text_frame.paragraphs[0]
p.text = "Graph-Native Agent Intelligence Layer"
p.font.size = Pt(28)
p.font.color.rgb = AMBER
p.font.name = "Calibri"

tb = s.shapes.add_textbox(Inches(0.6), Inches(4.7), Inches(12.4), Inches(0.6))
p = tb.text_frame.paragraphs[0]
p.text = ("Compile-time intelligence. Detection-only governance. "
          "Audit-first observability.")
p.font.size = Pt(18)
p.font.italic = True
p.font.color.rgb = LIGHT
p.font.name = "Calibri"

tb = s.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(12.4), Inches(0.5))
p = tb.text_frame.paragraphs[0]
p.text = "Executive Summary  ·  v0.1  ·  Phases 1–7 shipped"
p.font.size = Pt(14)
p.font.color.rgb = LIGHT
p.font.name = "Calibri"


# ---- 2. The problem -----------------------------------------------------
s = new()
add_title(s, "The problem we set out to solve")
add_subtitle(s, "Every LLM agent breaks the same way, for the same reasons.")
add_bullets(s, [
    ("Silent tool drift.",
     "An API version changes; the agent's tool binding is stale; the LLM keeps calling the old endpoint until users complain."),
    ("No cited memory.",
     "The agent has no way to say 'I did X because we learned Y last week.' Every decision is unmoored from evidence."),
    ("LLM at runtime for governance.",
     "Teams bolt a second LLM in front to 'check' the first — cost doubles, latency spikes, and the checker itself drifts."),
    ("No audit trail regulators trust.",
     "For the EU AI Act, SOC-2, or an internal incident review, 'the LLM decided' is not an acceptable answer."),
    ("Agents don't learn from lived experience.",
     "Failures repeat. What worked yesterday isn't remembered as a rule; it lives in a slack thread."),
])


# ---- 3. Why we built it -------------------------------------------------
s = new()
add_title(s, "Why Graxella")
add_subtitle(s, "A thin, graph-native layer that makes the fixes cheap AND cited.")
add_bullets(s, [
    ("Compile-time > runtime.",
     "Rules are promoted by humans and evaluated deterministically at dispatch. No runtime LLM for routing decisions."),
    ("Detection-only governance.",
     "Graxella observes and cites; it doesn't quietly rewrite in production. Humans approve; the runtime obeys."),
    ("Audit-first.",
     "Every Episode, Proposal, Promotion, and Rule is a node in a PROV-O JSON-LD graph — regulator-ready by construction."),
    ("Intelligence BEFORE traffic.",
     "Docs already describe most drift. Mine them; propose rules; let a reviewer promote before the first user hits the failure."),
    ("Same LLM, same tools — one line to wrap.",
     "graxella.wrap(app, tools=..., store=..., rulebook=...) is the entire integration."),
])


# ---- 4. What we built (architecture) -----------------------------------
s = new()
add_title(s, "What we built  —  architecture at a glance")
add_subtitle(s, "Five surfaces, one integrated python package: `graxella/`")

# top row: run-time surfaces
add_box(s, 0.6, 1.9, 3.0, 1.1, fill=AMBER, text_colour=NAVY, bold=True,
        text="wrap() / wrap_tools()\nlangchain + LangGraph shim", size=13)
add_box(s, 3.9, 1.9, 3.0, 1.1, fill=AMBER, text_colour=NAVY, bold=True,
        text="heal.route()\nrulebook-aware dispatch", size=13)
add_box(s, 7.2, 1.9, 3.0, 1.1, fill=AMBER, text_colour=NAVY, bold=True,
        text="GraxellaCallback\nEpisode capture", size=13)
add_box(s, 10.5, 1.9, 2.4, 1.1, fill=AMBER, text_colour=NAVY, bold=True,
        text="MCP server\nreview surface", size=13)

# middle row: brain
add_box(s, 0.6, 3.2, 6.0, 1.1, fill=WHITE, text_colour=NAVY, bold=True,
        text="Rulebook  (hot-reloadable JSON, atomic write, deterministic ids)", size=14)
add_box(s, 6.9, 3.2, 6.0, 1.1, fill=WHITE, text_colour=NAVY, bold=True,
        text="Hidden Agenda  (RuleDistiller · CapabilityReweigher · TrustPromoter · DocsMiner · DatalogMiner)", size=12)

# bottom row: memory + audit
add_box(s, 0.6, 4.5, 4.0, 1.1, fill=LIGHT, text_colour=NAVY,
        text="SqliteExperienceStore\ndurable Episode ledger", size=13)
add_box(s, 4.9, 4.5, 4.0, 1.1, fill=LIGHT, text_colour=NAVY,
        text="KnowledgeSeed  +  from_docs()\ntyped assertions from markdown", size=13)
add_box(s, 9.2, 4.5, 3.7, 1.1, fill=LIGHT, text_colour=NAVY,
        text="PROV-O audit export\nJSON-LD graph", size=13)

# arrows from run-time down into rulebook
add_arrow(s, 2.1, 3.0, 2.1, 3.2)
add_arrow(s, 5.4, 3.0, 5.4, 3.2)
add_arrow(s, 8.7, 3.0, 8.7, 3.2)

# arrows from middle into memory / audit
add_arrow(s, 3.6, 4.3, 2.6, 4.5)
add_arrow(s, 9.9, 4.3, 6.9, 4.5)
add_arrow(s, 9.9, 4.3, 11.05, 4.5)

# caption
tb = s.shapes.add_textbox(Inches(0.6), Inches(5.9), Inches(12.4), Inches(1.0))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Nothing is auto-applied. Miners emit Proposals; a human "
          "promotes; the runtime consults the Rulebook. PROV-O covers "
          "every hop from evidence to approval to dispatch.")
p.font.size = Pt(13); p.font.italic = True
p.font.color.rgb = LIGHT; p.font.name = "Calibri"


# ---- 5. Technical components (deep) -------------------------------------
s = new()
add_title(s, "Technical components  —  what each file does")
add_bullets(s, [
    ("graxella.experience / .persistence",
     "ExperienceEpisode dataclass · InMemory + SqliteExperienceStore. One Episode per node run, capturing tool_calls, latency, ok/err, decisions."),
    ("graxella.langgraph.callback",
     "GraxellaCallback — a plain BaseCallbackHandler; drops into any LangGraph or LangChain runnable via config={'callbacks':[…]}"),
    ("graxella.langgraph.healing",
     "TransformRecipe (field_map · static_defaults · drop_fields), HealedTool, heal.route() — the rulebook-aware dispatch primitive."),
    ("graxella.rulebook",
     "JSON on disk, atomic write, hot-reload, deterministic proposal ids (sha256(kind|subject)) so audit refs stay closed across re-mines."),
    ("graxella.hidden_agenda",
     "5 miners: RuleDistiller · CapabilityReweigher · TrustPromoter · DocsMiner · DatalogMiner (Phase 7). All emit the same Proposal shape."),
    ("graxella.ingest",
     "from_docs() → KnowledgeSeed. Five regex extractors turn markdown into typed Assertions (deprecated_by, field_maps_to, has_intent, has_arg, described_as)."),
    ("graxella.rules  (Phase 7)",
     "Minimal bottom-up Datalog engine (semi-naive fixpoint, safety-rule check) + facts bridge. Derives transitive substitutions from single-hop docs."),
    ("graxella.audit",
     "export() → W3C PROV-O JSON-LD. Every Episode/Miner/Proposal/Reviewer/Promotion/Rule is a graph node; cross-refs resolve."),
    ("graxella.mcp",
     "7 pure-Python handlers + optional stdio MCP server. Any MCP-aware client (Claude Desktop, Cursor) drives the review loop."),
    ("graxella.wrap  +  graxella.cli",
     "wrap()/wrap_tools() = the one-line integration. `graxella agenda run|review|promote|reject|rules|audit` = the terminal review loop."),
], size=13, spacing=4)


# ---- 6. ADLC page 1 — Design & Build ------------------------------------
s = new()
add_title(s, "ADLC  —  Graxella at every stage  (1/2)")
add_subtitle(s, "Design → Build → Test  ·  intelligence before traffic")

col1_x, col2_x, col3_x = 0.6, 4.7, 8.8
col_w = 4.0

add_box(s, col1_x, 1.9, col_w, 0.6, fill=AMBER, text_colour=NAVY, bold=True,
        text="1. DESIGN", size=16)
add_box(s, col1_x, 2.6, col_w, 4.0, fill=WHITE, text_colour=NAVY, align=PP_ALIGN.LEFT,
        text=(
            "•  Author API docs in markdown\n"
            "•  Deprecations expressed as prose or migration tables\n"
            "•  Intent labels (weather_lookup, notify_user…) declared\n\n"
            "Graxella action:\n"
            "•  ingest.from_docs() lifts typed Assertions from these docs\n"
            "•  Predicate vocabulary is small and stable — no NLP, no LLM"
        ), size=12)

add_box(s, col2_x, 1.9, col_w, 0.6, fill=AMBER, text_colour=NAVY, bold=True,
        text="2. BUILD", size=16)
add_box(s, col2_x, 2.6, col_w, 4.0, fill=WHITE, text_colour=NAVY, align=PP_ALIGN.LEFT,
        text=(
            "•  Wrap the LangGraph/LangChain app in one line:\n"
            "   graxella.wrap(app, tools=…, store=…,\n"
            "                 rulebook=…, docs=…)\n\n"
            "Graxella action:\n"
            "•  DocsMiner emits pre-traffic Proposals\n"
            "•  DatalogMiner derives multi-hop / transitive rules\n"
            "•  Reviewer promotes before first user hits the failure"
        ), size=12)

add_box(s, col3_x, 1.9, col_w, 0.6, fill=AMBER, text_colour=NAVY, bold=True,
        text="3. TEST", size=16)
add_box(s, col3_x, 2.6, col_w, 4.0, fill=WHITE, text_colour=NAVY, align=PP_ALIGN.LEFT,
        text=(
            "•  Run the wrapped app against a test store\n"
            "•  Assert Episode shape (tools=[…] shows the new path)\n\n"
            "Graxella action:\n"
            "•  Every test call becomes an Episode row\n"
            "•  PROV-O export gives a diff-able audit bundle\n"
            "•  heal.route() bypass is deterministic → no LLM flakiness"
        ), size=12)


# ---- 7. ADLC page 2 — Deploy / Operate / Learn --------------------------
s = new()
add_title(s, "ADLC  —  Graxella at every stage  (2/2)")
add_subtitle(s, "Deploy → Operate → Learn  ·  drift closes the loop back to Design")

col1_x, col2_x, col3_x = 0.6, 4.7, 8.8
col_w = 4.0

add_box(s, col1_x, 1.9, col_w, 0.6, fill=AMBER, text_colour=NAVY, bold=True,
        text="4. DEPLOY", size=16)
add_box(s, col1_x, 2.6, col_w, 4.0, fill=WHITE, text_colour=NAVY, align=PP_ALIGN.LEFT,
        text=(
            "•  Ship the wrapped app + rulebook.json + store.db\n"
            "•  Rulebook is hot-reloadable — new promotions land\n"
            "   on the next dispatch, no restart\n\n"
            "Graxella action:\n"
            "•  Atomic write on promote (.tmp → replace)\n"
            "•  Sub-ms dispatch overhead: JSON lookup + field_map"
        ), size=12)

add_box(s, col2_x, 1.9, col_w, 0.6, fill=AMBER, text_colour=NAVY, bold=True,
        text="5. OPERATE", size=16)
add_box(s, col2_x, 2.6, col_w, 4.0, fill=WHITE, text_colour=NAVY, align=PP_ALIGN.LEFT,
        text=(
            "•  Live traffic hits the wrapped app\n"
            "•  Reviewer opens MCP-connected Claude Desktop / Cursor\n"
            "   and drives list_proposals / promote / audit\n\n"
            "Graxella action:\n"
            "•  Each run → Episode → SqliteExperienceStore\n"
            "•  export_audit() emits PROV-O bundle on demand"
        ), size=12)

add_box(s, col3_x, 1.9, col_w, 0.6, fill=AMBER, text_colour=NAVY, bold=True,
        text="6. LEARN", size=16)
add_box(s, col3_x, 2.6, col_w, 4.0, fill=WHITE, text_colour=NAVY, align=PP_ALIGN.LEFT,
        text=(
            "•  Nightly HiddenAgendaRunner scans the store\n"
            "•  RuleDistiller finds (fail→success) pairs\n"
            "•  CapabilityReweigher / TrustPromoter refine routing\n\n"
            "Graxella action:\n"
            "•  New Proposals surface in the MCP review queue\n"
            "•  Loop feeds back to DESIGN — docs get updated to\n"
            "   match what the runtime already learned"
        ), size=12)


# ---- 8. Value proposition — for engineering ---------------------------
s = new()
add_title(s, "Value proposition  —  for engineering leaders")
add_bullets(s, [
    ("One-line integration.",
     "graxella.wrap(app, tools=..., store=..., rulebook=...) — no rewrite of node code, no framework switch."),
    ("No runtime LLM in the critical path.",
     "Routing is a rulebook lookup + field-map rename. Sub-millisecond, deterministic, replayable."),
    ("Cost predictability.",
     "Every OSS-model demo in this build runs locally on Ollama (qwen2.5:3b, 2GB). Cloud LLM spend is opt-in, not baseline."),
    ("Framework-friendly.",
     "Works with LangChain, LangGraph today; CrewAI + AutoGen adapters queued (Phase 8). Bring-your-own transport."),
    ("Small surface, verifiable code.",
     "Whole runtime is < 3000 LOC of pure python, zero heavy deps. Datalog engine is ~200 LOC and unit-testable."),
    ("Failure isolation.",
     "A stale tool binding is a rule promotion, not a hotfix release. Blast radius = one JSON entry."),
])


# ---- 9. Value proposition — for governance -----------------------------
s = new()
add_title(s, "Value proposition  —  for governance, risk & compliance")
add_bullets(s, [
    ("Detection-only by design.",
     "Graxella never auto-applies. Every rule change carries a human signature (approved_by) with a timestamp."),
    ("EU AI Act, SOC-2 ready.",
     "PROV-O JSON-LD audit graph covers Entity/Activity/Agent — the same shape auditors already know from W3C PROV."),
    ("Complete evidence chain.",
     "Every ApprovedRule cites derived_from — either Episode ids (lived) or doc source_uris (documented) — resolvable to the exact line."),
    ("Reproducibility guaranteed.",
     "Proposal ids are deterministic (sha256(kind|subject)). Re-mining the same store yields the same ids — audits don't dangle."),
    ("Human-in-the-loop is the default.",
     "The runtime reads the Rulebook; the Rulebook only grows via .promote(). No back-channel."),
    ("Reversible.",
     "Rules can be rejected; the audit records both the promotion AND the rejection. Full history is on disk, diff-able."),
])


# ---- 10. Value proposition — for the business --------------------------
s = new()
add_title(s, "Value proposition  —  for the business")
add_bullets(s, [
    ("Faster mean-time-to-recovery on tool drift.",
     "Documented deprecation → mined proposal → human promotion → live bypass. Minutes, not a hotfix cycle."),
    ("Institutional memory that survives team turnover.",
     "Every 'we learned last quarter that X' becomes a cited rule the next engineer inherits by default."),
    ("Vendor & model portability.",
     "The rulebook is model-agnostic. Swap qwen2.5:3b for GPT-4o or Claude Sonnet without touching a rule."),
    ("Lower LLM bill.",
     "Compile-time decisions displace runtime LLM calls that would otherwise be inside every request path."),
    ("Regulator-ready without a bolt-on.",
     "Audit isn't an integration; it's how Graxella stores state to begin with."),
    ("Foundation for the AXON program.",
     "Same primitives will host B8 (semantic tool discovery), B10 (payload healing), B11 (A2A routing), B12 (memory)."),
])


# ---- 11. Roadmap & status -----------------------------------------------
s = new()
add_title(s, "Where we are  &  what's next")

# progress table
def row(x, y, label, status, colour):
    add_box(s, x, y, 4.0, 0.5, fill=WHITE, text_colour=NAVY, align=PP_ALIGN.LEFT,
            text="   " + label, size=13, bold=True)
    add_box(s, x + 4.05, y, 1.6, 0.5, fill=colour, text_colour=WHITE,
            align=PP_ALIGN.CENTER, text=status, size=12, bold=True)

y = 1.8
row(0.6, y + 0*0.6, "Phase 1  ·  Experience + Store + Callback", "DONE", GREEN)
row(0.6, y + 1*0.6, "Phase 2  ·  axon_runtime intercept adapter", "DONE", GREEN)
row(0.6, y + 2*0.6, "Phase 3  ·  Rulebook + heal.route + PROV-O", "DONE", GREEN)
row(0.6, y + 3*0.6, "Phase 4  ·  from_docs + DocsMiner", "DONE", GREEN)
row(0.6, y + 4*0.6, "Phase 5  ·  wrap() + router node", "DONE", GREEN)
row(0.6, y + 5*0.6, "Phase 6  ·  MCP handlers + stdio server", "DONE", GREEN)
row(0.6, y + 6*0.6, "Phase 7  ·  Datalog engine + DatalogMiner", "DONE", GREEN)
row(0.6, y + 7*0.6, "Phase 8  ·  CrewAI + AutoGen adapters", "NEXT", AMBER)

# right column: highlights of the last shipped phase
add_box(s, 7.0, 1.8, 5.9, 4.9, fill=WHITE, text_colour=NAVY, align=PP_ALIGN.LEFT,
        text=(
            "Latest ship: Phase 7  (Datalog rule backend)\n\n"
            "•  200-LOC bottom-up Datalog engine (safe rules, semi-naive)\n"
            "•  Derives transitive substitutions single-hop docs miss\n"
            "•  Composes field_maps along multi-hop chains automatically\n"
            "•  Verified end-to-end against real OSS LLM\n"
            "   (Ollama qwen2.5:3b, ~2GB, runs local)\n\n"
            "Demo proved: LLM bound only to the deprecated get_weather,\n"
            "Graxella transparently reroutes to fetch_forecast_v3 with\n"
            "city→location rename — LLM never sees the failure.\n\n"
            "Next: adapters for CrewAI + AutoGen, then hardening (Phase 9)."
        ), size=12)


# ---- 12. Closing --------------------------------------------------------
s = new()
add_title(s, "Bottom line", y=2.0, size=44)

tb = s.shapes.add_textbox(Inches(0.6), Inches(3.2), Inches(12.4), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True

lines = [
    ("One line to wrap.",
     "Every agent call becomes cited, auditable, replayable."),
    ("Zero runtime LLM for routing.",
     "Compile-time rules, human-approved, sub-ms dispatch."),
    ("Regulator-ready from day one.",
     "PROV-O is the storage format, not a report generator."),
    ("Foundation for AXON  —  the Kubernetes for AI Agents.",
     ""),
]
for i, (h, b) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = ""
    r = p.add_run(); r.text = h
    r.font.size = Pt(24); r.font.bold = True
    r.font.color.rgb = AMBER; r.font.name = "Calibri"
    if b:
        r2 = p.add_run(); r2.text = "  " + b
        r2.font.size = Pt(20); r2.font.color.rgb = WHITE
        r2.font.name = "Calibri"
    p.space_after = Pt(10)

tb = s.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(12.4), Inches(0.5))
p = tb.text_frame.paragraphs[0]
p.text = "Contact:  ram@graxella  ·  ramprasad.vce@gmail.com"
p.font.size = Pt(14); p.font.italic = True
p.font.color.rgb = LIGHT; p.font.name = "Calibri"


# ---- write ---------------------------------------------------------------
OUT = Path(__file__).parent / "Graxella_Executive_Summary.pptx"
prs.save(OUT)
print(f"wrote {OUT}  ({OUT.stat().st_size / 1024:.1f} KB, {page[0]} slides)")
