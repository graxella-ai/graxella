"""Generate the Graxella flagship deck — "The Evidence Loop".

Run:  python build_flagship_ppt.py
Output: Graxella_Flagship_Deck.pptx (in this folder).

Covers: the enterprise problem (3 slides), current-stack vs graxella
direction, why now, the Evidence Loop thesis + diagram, production
failures vs graxella mechanisms, core technical values, architecture,
mnema + agent2society + gate deep dives, seamless adoption path,
roadmap, value metrics, moat.

Uses python-pptx (>=1.0). All colours + fonts set explicitly so the deck
renders identically regardless of the machine's PowerPoint theme.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- palette (evidence-ledger identity: paper / ink / verdigris / amber) ---
PAPER = RGBColor(0xFA, 0xFB, 0xF9)
PANEL = RGBColor(0xF0, 0xF4, 0xF0)
INK = RGBColor(0x1C, 0x25, 0x21)
MUTED = RGBColor(0x5C, 0x6A, 0x63)
GREEN = RGBColor(0x14, 0x80, 0x5E)   # verdigris accent — the ledger
AMBER = RGBColor(0xA8, 0x69, 0x1C)   # enterprise pain markers
RED = RGBColor(0xA9, 0x43, 0x2E)     # failure column
LINE = RGBColor(0xDC, 0xE3, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_TINT = RGBColor(0xE2, 0xEF, 0xE9)

HEAD_FONT = "Segoe UI"
BODY_FONT = "Segoe UI"
MONO_FONT = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ------------------------------------------------------------------ helpers
def new_slide():
    return prs.slides.add_slide(BLANK)


def add_bg(slide, colour=PAPER):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = colour
    bg.shadow.inherit = False
    return bg


def _set(p, text, *, size, colour, font=BODY_FONT, bold=False,
         italic=False, align=None, spacing=None):
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = colour
    p.font.name = font
    if align is not None:
        p.alignment = align
    if spacing is not None:
        p.space_after = Pt(spacing)
    return p


def add_eyebrow(slide, text, *, y=0.42, colour=GREEN, x=0.7, w=12.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.32))
    _set(tb.text_frame.paragraphs[0], text.upper(), size=12, colour=colour,
         font=MONO_FONT, bold=True)
    return tb


def add_title(slide, text, *, y=0.72, size=30, colour=INK, x=0.7, w=12.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.85))
    _set(tb.text_frame.paragraphs[0], text, size=size, colour=colour,
         font=HEAD_FONT, bold=True)
    return tb


def add_sub(slide, text, *, y=1.42, size=14, colour=MUTED, x=0.7, w=12.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    _set(tf.paragraphs[0], text, size=size, colour=colour)
    return tb


def add_bullets(slide, items, *, x=0.7, y=2.05, w=11.9, h=4.9,
                size=14, spacing=10, head_colour=GREEN, body_colour=INK):
    """items: list of str or (head, body) tuples."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        if isinstance(item, tuple):
            head, body = item
            run = p.add_run()
            run.text = f"{head} — "
            run.font.bold = True
            run.font.size = Pt(size)
            run.font.color.rgb = head_colour
            run.font.name = HEAD_FONT
            run2 = p.add_run()
            run2.text = body
            run2.font.size = Pt(size)
            run2.font.color.rgb = body_colour
            run2.font.name = BODY_FONT
        else:
            _set(p, item, size=size, colour=body_colour)
    return tb


def add_panel(slide, x, y, w, h, *, fill=PANEL, line=LINE, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    if radius:
        try:
            sh.adjustments[0] = 0.06
        except Exception:
            pass
    return sh


def panel_text(sh, lines, *, size=12, colour=INK, font=MONO_FONT,
               align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=4,
               margin=0.14):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin * 0.7)
    tf.margin_bottom = Inches(margin * 0.7)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, tuple):
            text, colr, bold = line
            _set(p, text, size=size, colour=colr, font=font, bold=bold,
                 align=align, spacing=spacing)
        else:
            _set(p, line, size=size, colour=colour, font=font,
                 align=align, spacing=spacing)
    return sh


def add_footer(slide, text="GRAXELLA · THE EVIDENCE LOOP"):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(7.08),
                                  Inches(12.0), Inches(0.3))
    _set(tb.text_frame.paragraphs[0], text, size=9, colour=MUTED,
         font=MONO_FONT)
    return tb


def add_table(slide, rows, *, x=0.7, y=2.0, w=11.9, col_widths=None,
              row_h=0.52, header_h=0.42, size=12, header_colour=GREEN):
    n_rows, n_cols = len(rows), len(rows[0])
    total_h = header_h + row_h * (n_rows - 1)
    gfx = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                 Inches(w), Inches(total_h))
    table = gfx.table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            table.columns[ci].width = Inches(cw)
    table.rows[0].height = Inches(header_h)
    for ri in range(1, n_rows):
        table.rows[ri].height = Inches(row_h)
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (GREEN_TINT if ri == 0 else
                                        (PAPER if ri % 2 else PANEL))
            cell.margin_left = Inches(0.09)
            cell.margin_right = Inches(0.09)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if ri == 0:
                _set(p, cell_text.upper(), size=size - 1, colour=header_colour,
                     font=MONO_FONT, bold=True)
            else:
                bold = ci == 0
                _set(p, cell_text, size=size,
                     colour=INK if ci == 0 else INK, bold=bold)
    return gfx


# =========================================================== 1 · TITLE
s = new_slide()
add_bg(s, INK)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.4))
_set(tb.text_frame.paragraphs[0], "GRAXELLA · FLAGSHIP PROGRAM", size=14,
     colour=RGBColor(0x43, 0xC4, 0x95), font=MONO_FONT, bold=True)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.3))
_set(tb.text_frame.paragraphs[0], "The Evidence Loop", size=54, colour=WHITE,
     font=HEAD_FONT, bold=True)
tb = s.shapes.add_textbox(Inches(0.9), Inches(3.8), Inches(11.0), Inches(0.9))
tf = tb.text_frame
tf.word_wrap = True
_set(tf.paragraphs[0],
     "Reliable AI agents at enterprise scale — the runtime where reliability "
     "compounds with volume, and every behavior change ships with citations.",
     size=19, colour=RGBColor(0xC9, 0xD4, 0xCE))
tb = s.shapes.add_textbox(Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.4))
_set(tb.text_frame.paragraphs[0],
     "DRAFT v1  ·  2026-08  ·  GRAXELLA + MNEMA + AGENT2SOCIETY + AXON + BROWNBRILLION",
     size=10, colour=RGBColor(0x7E, 0x8D, 0x85), font=MONO_FONT)

# ================================================ 2 · PROBLEM 1/3 — the math
s = new_slide()
add_bg(s)
add_eyebrow(s, "The problem · 1 of 3")
add_title(s, "Pilot purgatory: the math kills demos in production")
add_sub(s, "Agent chains multiply per-step error. The demo is one step; the "
           "business process is ten.")
stats = [("99%", "per-step", "≈ 90%", "over a 10-step chain"),
         ("95%", "per-step", "≈ 60%", "over a 10-step chain"),
         ("90%", "per-step", "≈ 35%", "over a 10-step chain")]
for i, (a, al, b, bl) in enumerate(stats):
    px = 0.7 + i * 4.1
    p = add_panel(s, px, 2.15, 3.8, 1.9, radius=True)
    tf = p.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(tf.paragraphs[0], f"{a} {al}", size=16, colour=MUTED, font=MONO_FONT,
         align=PP_ALIGN.CENTER)
    p2 = tf.add_paragraph()
    _set(p2, b, size=34, colour=RED if i else INK, font=HEAD_FONT, bold=True,
         align=PP_ALIGN.CENTER)
    p3 = tf.add_paragraph()
    _set(p3, bl, size=12, colour=MUTED, align=PP_ALIGN.CENTER)
add_bullets(s, [
    ("The demo works", "one agent, three steps, a friendly input. The production "
     "workflow is ten steps across systems that drift weekly."),
    ("Nothing improves with usage", "today's stacks have no mechanism that turns "
     "volume into reliability. Month six is exactly as fragile as week one."),
    ("So programs stall at pilot", "the industry's answer — retry with a smarter "
     "prompt — raises cost, not the curve."),
], y=4.45, size=14)
add_footer(s)

# ============================================ 3 · PROBLEM 2/3 — invisible value
s = new_slide()
add_bg(s)
add_eyebrow(s, "The problem · 2 of 3")
add_title(s, "Value is invisible — so budgets die at renewal")
add_sub(s, "Four compounding blockers keep enterprises from seeing (or proving) "
           "any return on agent programs.")
blocks = [
    ("INVISIBLE ROI", "Cost per completed task is opaque. Retries burn tokens "
     "silently. Nobody can show the CFO a curve moving the right way.", AMBER),
    ("THE AUDIT WALL", "EU AI Act, model risk management, SOX change control. "
     "“The LLM decided” is not an answer a regulator accepts.", AMBER),
    ("OPS DON'T COMPOUND", "Every incident ends with a human editing a prompt — "
     "unversioned, unreviewed, unmeasured. Headcount scales with agent count.", AMBER),
    ("FRAMEWORK FEAR", "LangGraph, CrewAI, AutoGen churn quarterly. No one wants "
     "to bet a program on an API surviving the year.", AMBER),
]
for i, (head, body, colr) in enumerate(blocks):
    px = 0.7 + (i % 2) * 6.15
    py = 2.1 + (i // 2) * 2.35
    p = add_panel(s, px, py, 5.85, 2.1, radius=True)
    tf = p.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    _set(tf.paragraphs[0], head, size=13, colour=colr, font=MONO_FONT, bold=True,
         spacing=6)
    _set(tf.add_paragraph(), body, size=13.5, colour=INK)
add_footer(s)

# ==================================== 4 · PROBLEM 3/3 — current vs graxella
s = new_slide()
add_bg(s)
add_eyebrow(s, "The problem · 3 of 3")
add_title(s, "Where the industry is vs where Graxella is heading")
add_table(s, [
    ["Dimension", "Current agent stacks", "Graxella direction"],
    ["Learning", "Runtime LLM judges, critics, retries — pay per call, forever",
     "Compile-time learning; deterministic O(1) dispatch"],
    ["Fixes", "Prompt edits in prose — unversioned, unreviewable",
     "Versioned promotions — gated, cited, rollbackable"],
    ["Audit", "Trace viewers; ‘why’ is forensics after the incident",
     "why() is a lookup — citations recorded by construction"],
    ["Governance", "Static guardrail rules or LLM-as-judge opinions",
     "Bayesian gate reading the org's own outcome ledger"],
    ["Reliability curve", "Flat — volume produces cost, not improvement",
     "Positive slope — volume produces evidence, evidence promotes fixes"],
    ["Lock-in", "Bet the program on one framework's API",
     "Substrate under any framework; adapters ride the churn"],
], y=1.95, col_widths=[2.0, 4.95, 4.95], row_h=0.72, size=12.5)
add_footer(s)

# ============================================================= 5 · WHY NOW
s = new_slide()
add_bg(s)
add_eyebrow(s, "Timing")
add_title(s, "Why now — four forces converge in 2026")
add_bullets(s, [
    ("The production wall is universal", "the first enterprise agent wave "
     "(2024–25) has hit scale-up meetings everywhere; the pain is budgeted, "
     "named, and unsolved."),
    ("Regulation has a date", "EU AI Act obligations are phasing in now. "
     "Cited, replayable decision trails move from nice-to-have to license-to-operate."),
    ("Token economics broke the LLM-judge pattern", "per-decision LLM governance "
     "gets linearly more expensive with volume — exactly when volume is the goal. "
     "A ledger lookup gets cheaper."),
    ("Enterprises already have the fuel", "millions of agent runs are being thrown "
     "away as logs. The evidence to make agents reliable already exists — nothing "
     "is built to learn from it."),
    ("The window", "observability vendors see traces, not decisions; framework "
     "vendors own syntax, not evidence. The evidence substrate seat is empty."),
], size=15, spacing=14, y=2.0)
add_footer(s)

# ============================================================ 6 · THE THESIS
s = new_slide()
add_bg(s)
add_eyebrow(s, "The solution approach")
add_title(s, "Agent behavior becomes an evidence-promoted artifact")
add_sub(s, "GitOps for agent behavior: the unit of improvement is a promotion, "
           "and the reviewer is a Bayesian gate reading your own operational history.")
add_bullets(s, [
    ("Invert the industry default", "everyone else adds runtime LLM judgment. "
     "Graxella moves ALL learning to compile time; the hot path stays "
     "deterministic, cheap, and auditable."),
    ("Every run produces evidence", "each dispatch auto-records a typed decision "
     "and a typed outcome into an immutable ledger (mnema). No opt-in, no "
     "instrumentation project."),
    ("Fixes are mined, not authored", "offline miners read the ledger and docs, "
     "and propose concrete changes: reroutes, field maps, healing recipes, skill "
     "reweights."),
    ("Promotion requires evidence", "a proposal ships only when the gate finds "
     "sufficient cited history for that exact (domain, kind) — or a human signs "
     "off. Constitution invariants are never bypassed."),
    ("Everything is reversible", "promotions are versioned with rollback points. "
     "Behavior change gets the change-management discipline code already has."),
], size=14.5, spacing=11, y=2.25)
add_footer(s)

# ================================================== 7 · EVIDENCE LOOP DIAGRAM
s = new_slide()
add_bg(s)
add_eyebrow(s, "The core mechanism")
add_title(s, "The Evidence Loop — six stages, one substrate")
stages = [
    ("DISPATCH", "deterministic · O(1)\nzero runtime LLM"),
    ("OUTCOME", "typed · auto-recorded\nper hop + per chain"),
    ("LEDGER", "immutable · cited\nmnema substrate"),
    ("PROPOSAL", "mined offline\nLLM allowed here"),
    ("GATE", "Bayesian · per-domain\nzero LLM in decision"),
    ("PROMOTION", "versioned\nrollbackable"),
]
bw, bh, gap = 1.92, 1.45, 0.12
x0, y0 = 0.72, 2.45
for i, (head, body) in enumerate(stages):
    bx = x0 + i * (bw + gap)
    box = add_panel(s, bx, y0, bw, bh, radius=True,
                    fill=GREEN_TINT if head == "LEDGER" else PANEL,
                    line=GREEN if head == "LEDGER" else LINE)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(tf.paragraphs[0], head, size=13, colour=GREEN if head == "LEDGER" else INK,
         font=MONO_FONT, bold=True, align=PP_ALIGN.CENTER, spacing=3)
    for j, ln in enumerate(body.split("\n")):
        _set(tf.add_paragraph(), ln, size=9.5, colour=MUTED,
             align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        ar = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(bx + bw - 0.02), Inches(y0 + bh / 2 - 0.09),
            Inches(gap + 0.06), Inches(0.18))
        ar.fill.solid()
        ar.fill.fore_color.rgb = GREEN
        ar.line.fill.background()
        ar.shadow.inherit = False
# return arrow
ret = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(x0 + 0.3),
                         Inches(y0 + bh + 0.55), Inches(11.4), Inches(0.42))
ret.fill.solid()
ret.fill.fore_color.rgb = GREEN_TINT
ret.line.color.rgb = GREEN
ret.line.width = Pt(1)
ret.shadow.inherit = False
tf = ret.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
_set(tf.paragraphs[0],
     "next dispatch reads promoted behavior — reliability compounds with volume",
     size=12, colour=GREEN, font=MONO_FONT, align=PP_ALIGN.CENTER)
add_bullets(s, [
    ("LLMs participate only offline", "mining proposals from evidence. The "
     "runtime path and the governance path are deterministic end to end."),
    ("One substrate, no side channels", "decisions, outcomes, proposals, gate "
     "verdicts, and promotions all live in the same immutable ledger — the audit "
     "trail is produced by the system, not bolted onto it."),
], y=5.35, size=13.5)
add_footer(s)

# ============================== 8 · PRODUCTION FAILURES VS GRAXELLA MECHANISM
s = new_slide()
add_bg(s)
add_eyebrow(s, "Failure handling")
add_title(s, "Production failures vs the Graxella mechanism")
add_table(s, [
    ["Production failure", "How it's handled today", "Graxella mechanism"],
    ["Tool deprecated / field renamed (drift)",
     "Agent burns retries until a human edits the prompt",
     "Drift mined from episodes → healing recipe (field map / reroute) gated, "
     "then applied at dispatch. Zero retries."],
    ["Chain breaks at step 7 of 10",
     "Whole workflow fails; root cause found by log archaeology",
     "Per-hop outcomes + trajectory ledger → chain-level healing proposal; same "
     "fault next week self-heals"],
    ["Task routed to the wrong agent",
     "Supervisor LLM guesses again (cost, non-determinism)",
     "Deterministic routing with audit features (matched tokens/tags); low-margin "
     "routes flagged for review"],
    ["Silent behavior drift after a change",
     "Discovered downstream, days later",
     "Detection-only governance: conflict, drift, low-confidence detectors fire "
     "loud events — nothing auto-corrects silently"],
    ["“Why did the agent do that?”",
     "Reconstructed by hand from traces, if at all",
     "why(decision) returns the citation chain — evidence, approver, threshold, "
     "rollback point"],
    ["A promoted fix regresses",
     "Roll back the prompt from memory, hope",
     "Promotions are versioned artifacts — one-step rollback, ledger records why"],
], y=1.9, col_widths=[2.85, 3.85, 5.2], row_h=0.83, size=11.5)
add_footer(s)

# ================================================ 9 · CORE TECHNICAL VALUES
s = new_slide()
add_bg(s)
add_eyebrow(s, "Core technical values")
add_title(s, "The four-pillar bar, held by four design principles")
pillars = [
    ("COMPLEXITY", "Scales to thousands of agents × hundreds of domains "
     "with zero new rules — evidence is per-(domain, kind), never global."),
    ("CLARITY", "Every number traceable to named evidence. Explanations are "
     "recorded at decision time, not reconstructed."),
    ("OPTIMIZATION", "Thresholds self-calibrate per tuple as outcomes accrue. "
     "Learning compounds locally, leaks nowhere."),
    ("RELIABILITY", "Deterministic hot path, O(1) governance, no LLM in any "
     "decision loop. Failures are loud, never swallowed."),
]
for i, (head, body) in enumerate(pillars):
    px = 0.7 + i * 3.08
    p = add_panel(s, px, 2.0, 2.9, 2.25, radius=True)
    tf = p.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.margin_top = Inches(0.14)
    _set(tf.paragraphs[0], head, size=13, colour=GREEN, font=MONO_FONT,
         bold=True, spacing=6)
    _set(tf.add_paragraph(), body, size=11.5, colour=INK)
add_bullets(s, [
    ("compile-time LLM > runtime LLM", "learning happens offline; dispatch is "
     "deterministic and cheap."),
    ("detection-only governance", "detectors flag, humans and evidence promote — "
     "nothing silently self-corrects, because silent remediation destroys audit."),
    ("memory-grounded, cited, self-calibrating", "every governance decision reads "
     "from the ledger and writes back to it."),
    ("silent plumbing, loud failures", "users write zero protocol/JSON/YAML; "
     "every degradation is a flagged, visible event."),
], y=4.55, size=13.5, spacing=8)
add_footer(s)

# ============================================================ 10 · ARCHITECTURE
s = new_slide()
add_bg(s)
add_eyebrow(s, "Architecture")
add_title(s, "Control plane / data plane — built for scale from the split")
dp = add_panel(s, 0.7, 2.0, 5.9, 4.35, radius=True)
tf = dp.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.16)
_set(tf.paragraphs[0], "DATA PLANE — in the customer's runtime", size=13,
     colour=GREEN, font=MONO_FONT, bold=True, spacing=8)
for txt in [
    "• graxella.mesh / instrument — one-line wrap over native agents",
    "• Deterministic router (agent2society) — TF-IDF or small local embed",
    "• Handoff runtime — bounded multi-hop, budgets, loop detection",
    "• Healing dispatch (brownbrillion) — promoted recipes applied at call time",
    "• Adapters: LangGraph · CrewAI · AutoGen · MCP · A2A · plain callables",
    "",
    "Embedded, framework-agnostic, no network dependency on the hot path.",
]:
    _set(tf.add_paragraph(), txt, size=12.5,
         colour=MUTED if txt.startswith("Embedded") else INK, spacing=6)
cp = add_panel(s, 6.85, 2.0, 5.75, 4.35, radius=True, fill=GREEN_TINT, line=GREEN)
tf = cp.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.16)
_set(tf.paragraphs[0], "CONTROL PLANE — multi-tenant service", size=13,
     colour=GREEN, font=MONO_FONT, bold=True, spacing=8)
for txt in [
    "• Mnema ledger — immutable evidence store (SQLite edge / Postgres scale)",
    "• Evidence Gate — Bayesian promotion decisions, O(1) per decision",
    "• Offline miners (axon + agenda) — episodes + docs → proposals",
    "• Audit service — W3C PROV-O export, change-management reports",
    "• Value ledger — ROI and reliability curves as queries",
    "",
    "Horizontally scalable, sharded by namespace · domain. Evidence never "
    "leaks across tenants or domains.",
]:
    _set(tf.add_paragraph(), txt, size=12.5,
         colour=MUTED if txt.startswith("Horizontally") else INK, spacing=6)
add_footer(s)

# ================================================================ 11 · MNEMA
s = new_slide()
add_bg(s)
add_eyebrow(s, "Component deep dive · 1 of 3")
add_title(s, "Mnema — the memory intelligence substrate")
add_sub(s, "Not a vector cache. A belief system with version control — the "
           "single source of truth every governance decision reads from.")
add_bullets(s, [
    ("Immutable typed assertions", "beliefs are never edited; revision = new "
     "assertion + supersedes link. History cannot be rewritten — by construction."),
    ("Bi-temporal by design", "valid_from/valid_to (true in the world) vs "
     "asserted_at (when learned). “What did we believe on May 3rd?” is a query."),
    ("Provenance is mandatory", "no anonymous beliefs. Every assertion carries "
     "origin (observed / told-by-user / told-by-agent / inferred / doc) and source id."),
    ("Content-addressed identity", "the hash covers propositional content only — "
     "confidence and status can evolve without forging identity."),
    ("Retraction cascade", "retract one belief and every belief derived from it "
     "is found and re-examined — the trust graph stays consistent."),
    ("Sleep/wake consolidation", "offline compression of episodes into stable "
     "beliefs (LLM allowed — it's compile time), with full citation of sources."),
    ("why() / timeline()", "any belief explains itself: who said it, what it "
     "derived from, what superseded it, when."),
], size=13.5, spacing=8, y=2.15)
add_footer(s)

# ========================================================= 12 · AGENT2SOCIETY
s = new_slide()
add_bg(s)
add_eyebrow(s, "Component deep dive · 2 of 3")
add_title(s, "agent2society — the deterministic agent-to-agent layer")
add_sub(s, "Routing, handoffs, and governance between agents — with zero "
           "runtime LLM and a written explanation for every decision.")
add_bullets(s, [
    ("Deterministic routing", "one embedding lookup (TF-IDF default, pluggable "
     "local models) scores every (agent, skill) pair — no supervisor LLM haggling "
     "in natural language."),
    ("Audit features on every score", "each candidate carries the matched tokens "
     "and overlapping tags that produced its number — a human can see why the "
     "score is what it is."),
    ("Handoff envelope", "task, intent, assumptions, confidence-required — a "
     "typed contract between agents, never user-authored JSON."),
    ("Conformance checking", "the top candidate is verified against the declared "
     "capability card; failures fall through to the next candidate, flagged."),
    ("Detection-only governance", "conflict, capability-drift, low-confidence, "
     "low-margin detectors fire structured events. Nothing auto-corrects silently."),
    ("Explanation store", "every routing decision is persisted as a replayable "
     "RoutingExplanation — the routing audit log doubles as training evidence."),
    ("Backtest-gated optimizer", "proposed skill-tag improvements must win on "
     "recorded history before they're even proposed for promotion."),
], size=13.5, spacing=8, y=2.15)
add_footer(s)

# ========================================================= 13 · EVIDENCE GATE
s = new_slide()
add_bg(s)
add_eyebrow(s, "Component deep dive · 3 of 3")
add_title(s, "The Evidence Gate — governance without rules or judges")
add_sub(s, "No scoring rubric, no rule DSL, no LLM judge. One question, answered "
           "from the ledger: what happened the last N times this change was made "
           "in this domain?")
p = add_panel(s, 0.7, 2.3, 7.1, 3.4, fill=PANEL, line=LINE, radius=True)
panel_text(p, [
    ("# gate.why(promotion_2841)", MUTED, False),
    ("decision:   AUTO_APPROVE", INK, True),
    ("posterior:  0.94  (Beta(47+1, 2+1), domain=refunds)", INK, False),
    ("evidence:   47 successes · 2 failures · 9 sessions", INK, False),
    ("citations:  [asr_e01f…, asr_a4c2…, asr_77b0…, +44]", GREEN, False),
    ("threshold:  0.90  (self-calibrated; cold-start=HUMAN)", INK, False),
    ("guards:     constitution PASS · blast=narrow · div=9≥3", INK, False),
], size=12.5, spacing=6)
add_bullets(s, [
    ("Bayesian, per-tuple", "closed-form posterior per (domain, kind, target). "
     "Cold domains: everything needs a human."),
    ("Self-calibrating", "thresholds loosen only for tuples with confirmed "
     "outcomes — and only there."),
    ("Poisoning defense", "loosening requires ≥K independent sessions and "
     "operators."),
    ("Hard blocks stay hard", "constitution invariants sit above the gate; "
     "evidence never bypasses them."),
    ("Self-auditing", "every verdict is itself a cited ledger assertion."),
], x=8.0, y=2.3, w=4.6, size=12, spacing=7)
add_footer(s)

# ======================================================= 14 · SEAMLESS SWITCH
s = new_slide()
add_bg(s)
add_eyebrow(s, "Adoption")
add_title(s, "Switching is one line — and reversible at every step")
p = add_panel(s, 0.7, 1.85, 11.9, 1.75, fill=PANEL, line=LINE, radius=True)
panel_text(p, [
    ("# your code today — unchanged", MUTED, False),
    ("triage    = create_react_agent(llm, [check_order, lookup_policy], name=\"triage\")", INK, False),
    ("responder = create_react_agent(llm, [draft_email], name=\"responder\")", INK, False),
    ("# the switch — one wrap, same call surface, nothing removed", MUTED, False),
    ("app = graxella.mesh([triage, responder], memory=graxella.Memory.sqlite(\"./mnema.db\"))", GREEN, True),
], size=12, spacing=4)
steps = [
    ("1 · OBSERVE", "Wrap existing agents. Framework code unchanged; graxella "
     "only records. Risk: zero."),
    ("2 · EVIDENCE", "Decisions + outcomes accrue in the ledger. Baseline "
     "cost/task and success dashboards appear."),
    ("3 · GOVERN", "Turn the gate on. Proposals flow to human review; audit "
     "answers become one-line queries."),
    ("4 · COMPOUND", "Warmed domains auto-promote with citations. Reliability "
     "slope goes positive."),
]
for i, (head, body) in enumerate(steps):
    px = 0.7 + i * 3.08
    pnl = add_panel(s, px, 3.95, 2.9, 1.85, radius=True)
    tf = pnl.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.12)
    _set(tf.paragraphs[0], head, size=12.5, colour=GREEN, font=MONO_FONT,
         bold=True, spacing=5)
    _set(tf.add_paragraph(), body, size=11.5, colour=INK)
add_bullets(s, [
    ("Keep your framework, keep both paths", "native LangGraph/CrewAI agents or "
     "graxella.Agent — both land in the same mesh; choice is always additive."),
    ("Runs alongside existing observability", "OTel traces delegate to your APM; "
     "graxella adds the decision layer APM can't see. Unwrap anytime — your "
     "agents still run."),
], y=6.0, size=12.5, spacing=6)
add_footer(s)

# ============================================================== 15 · ROADMAP
s = new_slide()
add_bg(s)
add_eyebrow(s, "Execution")
add_title(s, "Six phases, 26 weeks — every phase ends in buyer-visible value")
add_table(s, [
    ["Phase", "Weeks", "Delivers", "Exit proof"],
    ["0 · Truth flows", "1–3",
     "Auto-recorded typed outcomes; packaging, tests, loud failures",
     "Every routed task yields a cited decision+outcome pair"],
    ["1 · Evidence Gate", "3–8",
     "Bayesian gate replaces scored policy; one unified proposal pipeline",
     "Same proposal: blocked cold, auto-approved after N cited successes"],
    ["2 · Handoff Runtime", "8–14",
     "Bounded multi-hop dispatch; chain-level outcomes and healing",
     "Injected fault at hop 3 self-heals the following week"],
    ["3 · Scale substrate", "14–20",
     "Control/data-plane split; Postgres ledger; OTel; async",
     "Route p50 <10ms @ 1k agents; gate p50 <5ms; crash-safe"],
    ["4 · Value Ledger", "20–26",
     "ROI as query; reliability-slope dashboard; compliance pack",
     "“Wk 1: 71% @ $0.48/task → Wk 8: 94% @ $0.19” — all cited"],
    ["5 · Ecosystem", "26+",
     "CrewAI/AutoGen adapters; MCP healing; evidence federation",
     "Same workflow crosses frameworks; ledger and curve intact"],
], y=1.85, col_widths=[2.3, 0.95, 4.6, 4.05], row_h=0.78, size=11.5)
add_footer(s)

# ======================================================== 16 · VALUE VISIBLE
s = new_slide()
add_bg(s)
add_eyebrow(s, "Value made visible")
add_title(s, "The numbers an enterprise finally gets to see")
add_table(s, [
    ["Metric", "Definition", "Why it's the headline"],
    ["Reliability slope", "Δ task success per 1,000 runs",
     "The compounding claim, falsifiable. Everyone else's slope is ~zero."],
    ["Cost / completed task", "Spend ÷ successful completions, trended",
     "The CFO's number — falls as promotions accumulate."],
    ["MTTH", "Drift detected → healed in production",
     "Days (human) → minutes (gated auto-promotion)."],
    ["Auto-promotion rate", "% changes shipped by the gate, cited",
     "Governance scales without headcount."],
    ["Audit answerability", "% decisions with a complete why-chain",
     "100% by construction — the regulator metric."],
], y=1.9, col_widths=[2.5, 4.2, 5.2], row_h=0.62, size=12)
p = add_panel(s, 0.7, 5.5, 11.9, 1.15, fill=GREEN_TINT, line=GREEN, radius=True)
tf = p.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.25)
_set(tf.paragraphs[0],
     "The pilot report a VP reads without translation:  “Week 1: 71% success, "
     "$0.48/task.  Week 8: 94%, $0.19/task.  23 behavior changes shipped — every "
     "one cited, every one reversible.”",
     size=14, colour=GREEN, font=HEAD_FONT, bold=True)
add_footer(s)

# ================================================================ 17 · MOAT
s = new_slide()
add_bg(s)
add_eyebrow(s, "Defensibility")
add_title(s, "Why Graxella wins — and keeps winning")
add_bullets(s, [
    ("The corpus compounds", "every task run deepens the customer's own evidence "
     "ledger. Switching cost grows daily, honestly — the moat is the customer's "
     "history, which only graxella reads back as governance."),
    ("Auditability can't be retrofitted", "stacks built on runtime LLM judgment "
     "cannot recover citations they never recorded. Graxella's trail exists by "
     "construction, not by instrumentation."),
    ("The economics invert at scale", "per-decision LLM governance gets more "
     "expensive with volume; a ledger lookup gets cheaper. The bigger the "
     "deployment, the wider the advantage."),
    ("Framework churn is a tailwind", "the surface mirrors whatever the customer "
     "already writes; the value lives underneath. Every migration proves the "
     "evidence ledger outlives the stack."),
], size=15, spacing=14, y=2.0)
add_footer(s)

# =============================================================== 18 · CLOSE
s = new_slide()
add_bg(s, INK)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.5), Inches(0.4))
_set(tb.text_frame.paragraphs[0], "GRAXELLA · THE EVIDENCE LOOP", size=13,
     colour=RGBColor(0x43, 0xC4, 0x95), font=MONO_FONT, bold=True)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.8))
tf = tb.text_frame
tf.word_wrap = True
_set(tf.paragraphs[0],
     "Agents that get more reliable the more you run them —",
     size=34, colour=WHITE, font=HEAD_FONT, bold=True)
_set(tf.add_paragraph(),
     "with a citation for every change, and a rollback for every mistake.",
     size=34, colour=RGBColor(0x43, 0xC4, 0x95), font=HEAD_FONT, bold=True)
tb = s.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(11.3), Inches(1.4))
tf = tb.text_frame
tf.word_wrap = True
for i, line in enumerate([
    "compile-time LLM > runtime LLM   ·   detection-only governance",
    "memory-grounded, cited, self-calibrating   ·   silent plumbing, loud failures",
]):
    p_ = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    _set(p_, line, size=13, colour=RGBColor(0x9A, 0xA8, 0xA0), font=MONO_FONT,
         spacing=6)
tb = s.shapes.add_textbox(Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.4))
_set(tb.text_frame.paragraphs[0],
     "NEXT STEP: PHASE 0 — CLOSE THE LOOP (WEEKS 1–3)", size=11,
     colour=RGBColor(0x7E, 0x8D, 0x85), font=MONO_FONT)

# ------------------------------------------------------------------- save
out = Path(__file__).parent / "Graxella_Flagship_Deck.pptx"
prs.save(out)
print(f"Wrote {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
