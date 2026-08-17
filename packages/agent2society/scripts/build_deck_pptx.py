"""Generate a technical slide deck for agent2society as a .pptx file.

Widescreen 16:9, dark-on-light palette, code in monospace, deterministic
layouts (no auto-layout drift).
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = Path(__file__).resolve().parent.parent / "agent2society_technical_deck.pptx"

# --- Palette ----------------------------------------------------------------
INK = RGBColor(0x10, 0x14, 0x1F)        # primary text
MUTED = RGBColor(0x55, 0x5B, 0x6E)      # secondary text
ACCENT = RGBColor(0x2D, 0x5B, 0xFF)     # links / highlight
ACCENT_DARK = RGBColor(0x12, 0x2E, 0xB0)
GOOD = RGBColor(0x0E, 0x8C, 0x4F)
WARN = RGBColor(0xC8, 0x5A, 0x10)
BAD = RGBColor(0xB2, 0x1E, 0x35)
BG = RGBColor(0xFB, 0xFB, 0xFC)
BG_PANEL = RGBColor(0xF1, 0xF3, 0xF8)
BG_CODE = RGBColor(0x10, 0x14, 0x1F)
CODE_TEXT = RGBColor(0xEA, 0xEE, 0xF7)
RULE = RGBColor(0xC4, 0xCB, 0xD8)

# --- Fonts ------------------------------------------------------------------
HEAD = "Calibri"
BODY = "Calibri"
MONO = "Consolas"

# --- Slide geometry (16:9) --------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# --- Helpers ---------------------------------------------------------------
def add_bg(slide, color=BG):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.shadow.inherit = False
    return rect


def add_textbox(slide, left, top, width, height, *, fill=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.fill.background()
    return box, tf


def set_run(run, *, text=None, font=BODY, size=18, bold=False, italic=False, color=INK):
    if text is not None:
        run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_para(tf, text, *, font=BODY, size=18, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT, space_before=0, space_after=4, first=False, bullet=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if bullet:
        run = p.add_run()
        set_run(run, text="•  ", font=font, size=size, bold=False, color=MUTED)
        run2 = p.add_run()
        set_run(run2, text=text, font=font, size=size, bold=bold, italic=italic, color=color)
    else:
        run = p.add_run()
        set_run(run, text=text, font=font, size=size, bold=bold, italic=italic, color=color)
    return p


def add_title_bar(slide, title, eyebrow=None):
    """Top bar with optional eyebrow + slide title + thin rule."""
    if eyebrow:
        box, tf = add_textbox(slide, MARGIN, Inches(0.35), SLIDE_W - 2 * MARGIN, Inches(0.32))
        add_para(tf, eyebrow.upper(), font=HEAD, size=11, bold=True, color=ACCENT, first=True, space_after=0)
    box, tf = add_textbox(slide, MARGIN, Inches(0.62), SLIDE_W - 2 * MARGIN, Inches(0.7))
    add_para(tf, title, font=HEAD, size=30, bold=True, color=INK, first=True, space_after=0)
    # rule
    rule = slide.shapes.add_connector(1, MARGIN, Inches(1.4), SLIDE_W - MARGIN, Inches(1.4))
    rule.line.color.rgb = RULE
    rule.line.width = Pt(0.75)


def add_footer(slide, slide_num, total):
    box, tf = add_textbox(slide, MARGIN, SLIDE_H - Inches(0.42), SLIDE_W - 2 * MARGIN, Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    set_run(r, text="agent2society", font=HEAD, size=10, bold=True, color=MUTED)
    r2 = p.add_run()
    set_run(r2, text="   ·   v0.5.3   ·   pip install agent2society", font=HEAD, size=10, color=MUTED)
    # right-side number
    box2, tf2 = add_textbox(slide, SLIDE_W - Inches(1.2) - MARGIN, SLIDE_H - Inches(0.42), Inches(1.2), Inches(0.3))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r3 = p2.add_run()
    set_run(r3, text=f"{slide_num} / {total}", font=HEAD, size=10, color=MUTED)


def add_code_block(slide, code, left, top, width, height, *, size=12):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.04
    box.fill.solid()
    box.fill.fore_color.rgb = BG_CODE
    box.line.fill.background()
    box.shadow.inherit = False

    tb, tf = add_textbox(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), height - Inches(0.24))
    tf.word_wrap = False
    lines = code.split("\n")
    for i, line in enumerate(lines):
        add_para(tf, line if line else " ", font=MONO, size=size, color=CODE_TEXT,
                 first=(i == 0), space_after=0)


def add_panel(slide, left, top, width, height, *, fill=BG_PANEL):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.adjustments[0] = 0.04
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    panel.line.color.rgb = RULE
    panel.line.width = Pt(0.5)
    panel.shadow.inherit = False
    return panel


def add_pill(slide, left, top, text, *, fill=ACCENT, fg=RGBColor(0xFF, 0xFF, 0xFF), size=11):
    width = Inches(max(0.7, 0.16 * len(text) + 0.4))
    height = Inches(0.32)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    pill.shadow.inherit = False
    tf = pill.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(0)
    r = p.add_run()
    set_run(r, text=text, font=HEAD, size=size, bold=True, color=fg)
    return pill, width


# ============================================================================
# SLIDES
# ============================================================================

# Track total for footer numbering. We pass slide builders that return slide,
# then add footer at the end.
_slide_builders = []


def slide(fn):
    _slide_builders.append(fn)
    return fn


# ---------------------------------------------------------------------------
@slide
def s_title(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, BG)

    # Accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SLIDE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

    box, tf = add_textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.3))
    add_para(tf, "agent2society", font=HEAD, size=64, bold=True, color=INK, first=True, space_after=0)

    box, tf = add_textbox(s, Inches(0.9), Inches(3.2), Inches(11.5), Inches(1.5))
    add_para(tf,
             "The transparent coordination layer for A2A agent meshes.",
             font=HEAD, size=24, italic=True, color=MUTED, first=True, space_after=8)
    add_para(tf,
             "Deterministic routing. Conformance guardrails. Governance hooks. "
             "A human-readable explanation for every decision.",
             font=BODY, size=18, color=INK, space_after=0)

    # Pills
    add_pill(s, Inches(0.9), Inches(5.2), "v0.5.3", fill=ACCENT)
    add_pill(s, Inches(1.95), Inches(5.2), "Apache-2.0", fill=GOOD)
    add_pill(s, Inches(3.3), Inches(5.2), "Zero hard deps", fill=ACCENT_DARK)
    add_pill(s, Inches(5.05), Inches(5.2), "Python 3.9+", fill=MUTED)
    add_pill(s, Inches(6.3), Inches(5.2), "112 tests passing", fill=GOOD)

    box, tf = add_textbox(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.6))
    add_para(tf, "pip install agent2society", font=MONO, size=18, color=ACCENT, first=True, space_after=0)
    add_para(tf, "github.com/graxella/agent2society    ·    pypi.org/project/agent2society",
             font=HEAD, size=12, color=MUTED, space_after=0)


# ---------------------------------------------------------------------------
@slide
def s_problem(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Multi-agent systems are powerful and untrustworthy", eyebrow="The problem")

    # Two-column: left = the failure mode, right = root cause
    add_panel(s, MARGIN, Inches(1.7), Inches(6.0), Inches(5.2))
    box, tf = add_textbox(s, MARGIN + Inches(0.2), Inches(1.85), Inches(5.7), Inches(5.0))
    add_para(tf, "What you see in production", font=HEAD, size=14, bold=True, color=ACCENT, first=True, space_after=10)
    add_para(tf, "Finance question silently routed to the writer agent.", bullet=True, size=15, space_after=8)
    add_para(tf, "PII-laden ticket sent to a public-search agent.", bullet=True, size=15, space_after=8)
    add_para(tf, "Marketing memo cites a number nobody can trace.", bullet=True, size=15, space_after=8)
    add_para(tf, "Costs creep on every dispatch — supervisor reasons in every turn.", bullet=True, size=15, space_after=8)
    add_para(tf, "The transcript exists. The explanation does not.", bullet=True, size=15, bold=True, color=BAD, space_after=0)

    add_panel(s, Inches(6.85), Inches(1.7), Inches(6.0), Inches(5.2), fill=RGBColor(0xFB, 0xF1, 0xF1))
    box, tf = add_textbox(s, Inches(7.05), Inches(1.85), Inches(5.7), Inches(5.0))
    add_para(tf, "Root cause", font=HEAD, size=14, bold=True, color=BAD, first=True, space_after=10)
    add_para(tf,
             "The supervisor pattern hides every routing choice behind an LLM that \"decided\".",
             size=16, italic=True, color=INK, space_after=10)
    add_para(tf,
             "You inherit two costs:",
             size=15, bold=True, space_after=6)
    add_para(tf, "Tokens — the supervisor reasons in every turn.", bullet=True, size=15, space_after=6)
    add_para(tf, "Trust — routing is whatever the model felt like that turn.", bullet=True, size=15, space_after=10)
    add_para(tf,
             "Reflection / critic / explain-yourself agents push opacity one layer down — and bill you for the privilege.",
             size=14, italic=True, color=MUTED, space_after=0)


# ---------------------------------------------------------------------------
@slide
def s_thesis(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Replace the supervisor with a deterministic graph", eyebrow="The thesis")

    box, tf = add_textbox(s, MARGIN, Inches(1.8), SLIDE_W - 2 * MARGIN, Inches(1.5))
    add_para(tf,
             "Every routing decision is a deterministic graph traversal you can read, "
             "a score you can threshold on, and an audit trail keyed by handoff id.",
             font=HEAD, size=22, italic=True, color=INK, first=True, space_after=0)

    # Three pillars
    pillar_w = Inches(4.0)
    pillar_h = Inches(3.6)
    gap = Inches(0.2)
    pillars = [
        ("ROUTING + DISPATCH",
         "Capability graph parsed from A2A Agent Cards. Deterministic TF-IDF (or pluggable embedding) scorer. Graph-derived conformance guardrail. A2A JSON-RPC dispatcher.",
         ACCENT),
        ("MEANING + CONTEXT",
         "Handoff envelope carries intent, assumptions, and an upstream decision chain. SelfAssessment surfaces each agent's declared limits on every explanation that picks it.",
         ACCENT_DARK),
        ("TRANSPARENCY + GOVERNANCE",
         "RoutingExplanation for every decision (success OR failure). Detection-only hooks for low confidence, conflicts, capability drift, human review. Never silently auto-corrects.",
         GOOD),
    ]
    x = MARGIN
    for label, body, color in pillars:
        add_panel(s, x, Inches(3.5), pillar_w, pillar_h)
        # Color stripe at top
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.18), Inches(3.65), pillar_w - Inches(0.36), Inches(0.08))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = color; stripe.line.fill.background()
        box, tf = add_textbox(s, x + Inches(0.25), Inches(3.85), pillar_w - Inches(0.5), pillar_h - Inches(0.4))
        add_para(tf, label, font=HEAD, size=12, bold=True, color=color, first=True, space_after=10)
        add_para(tf, body, font=BODY, size=14, color=INK, space_after=0)
        x += pillar_w + gap


# ---------------------------------------------------------------------------
@slide
def s_cost(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "The cost wedge: one embedding instead of supervisor LLM rounds", eyebrow="Why now")

    # Big numbers row
    nums = [
        ("94.12%", "coordination tokens removed", GOOD),
        ("~140x", "cheaper coordination", ACCENT),
        ("11/12 vs 11/12", "parity on task success", INK),
        ("9 vs 427", "median tokens / task", ACCENT_DARK),
    ]
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.6)) / 4
    y = Inches(1.85)
    h = Inches(1.7)
    for i, (big, label, color) in enumerate(nums):
        x = MARGIN + i * (col_w + Inches(0.2))
        add_panel(s, x, y, col_w, h)
        box, tf = add_textbox(s, x + Inches(0.1), y + Inches(0.15), col_w - Inches(0.2), Inches(0.95))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, text=big, font=HEAD, size=34, bold=True, color=color)
        box, tf = add_textbox(s, x + Inches(0.1), y + Inches(1.1), col_w - Inches(0.2), Inches(0.5))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, text=label, font=HEAD, size=12, color=MUTED)

    # Benchmark code block
    code = (
        "coordination tokens per task\n"
        "  supervisor:     median=427  mean=427.8  min=424  max=433\n"
        "  agent2society:  median=9    mean=9.4    min=7    max=15\n"
        "\n"
        "reduction:       94.12%\n"
        "\n"
        "task success (right agent + skill chosen)\n"
        "  supervisor:     11/12\n"
        "  agent2society:  11/12"
    )
    add_code_block(s, code, MARGIN, Inches(3.85), Inches(7.5), Inches(2.85), size=13)

    # Notes column
    box, tf = add_textbox(s, Inches(8.4), Inches(3.85), Inches(4.4), Inches(2.9))
    add_para(tf, "Methodology", font=HEAD, size=14, bold=True, color=ACCENT, first=True, space_after=8)
    add_para(tf, "Dispatch cost (actual agent work) is held EQUAL across methods.",
             bullet=True, size=13, space_after=6)
    add_para(tf, "Delta is purely the supervisor's reasoning tokens.",
             bullet=True, size=13, space_after=6)
    add_para(tf, "Real LangGraph head-to-head: 96.11% reduction, ~200x ratio.",
             bullet=True, size=13, space_after=6)
    add_para(tf, "Parity, not perfection — claim verified by labelled suite.",
             bullet=True, size=13, space_after=0)


# ---------------------------------------------------------------------------
@slide
def s_explanation(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Every decision produces a human-readable explanation",
                  eyebrow="What an explanation looks like")

    code = (
        "handoff a77d6eed5a15  task: Draft an executive memo on customer churn\n"
        "  intent: prep the Q3 board pack\n"
        "  assumes:\n"
        "    - churn data is final\n"
        "  chose  : writer-agent :: exec_memo (confidence=0.593)\n"
        "  why    : writer-agent::exec_memo selected (score=0.593). Semantic match\n"
        "           0.725 on tokens [an, draft, executive, memo]; tag overlap 0.200\n"
        "           on [draft, memo]. Runner-up research-agent::web_research rejected:\n"
        "           below min_score.\n"
        "  features: min_score=0.050, confidence_required=0.500, score=0.593,\n"
        "            semantic=0.725, tag_overlap=0.200\n"
        "  alternatives:\n"
        "    - writer-agent :: exec_memo  score=0.593 sem=0.725 tags=0.200\n"
        "    - research-agent :: web_research  score=0.000  REJECTED (below min_score)\n"
        "  agent self-caveats:\n"
        "    - English only\n"
        "    - max ~400 words per memo\n"
        "    - escalate when: any quantitative claim cited without source"
    )
    add_code_block(s, code, MARGIN, Inches(1.7), Inches(8.3), Inches(5.2), size=11)

    box, tf = add_textbox(s, Inches(9.1), Inches(1.7), Inches(3.75), Inches(5.2))
    add_para(tf, "It's a template", font=HEAD, size=14, bold=True, color=ACCENT, first=True, space_after=8)
    add_para(tf, "No LLM rendered it. It's a deterministic render of features the router actually used.",
             size=13, space_after=12)
    add_para(tf, "ASCII-safe", font=HEAD, size=14, bold=True, color=ACCENT, space_after=8)
    add_para(tf, "Prints on a default Windows console — pinned by regression test.",
             size=13, space_after=12)
    add_para(tf, "Includes failures", font=HEAD, size=14, bold=True, color=ACCENT, space_after=8)
    add_para(tf, "Conformance blocks and dispatch errors get their own explanation with blocked_reason.",
             size=13, space_after=12)
    add_para(tf, "Keyed by handoff id", font=HEAD, size=14, bold=True, color=ACCENT, space_after=8)
    add_para(tf, "society.explain(h.id) pulls the rationale on any line in your ledger.",
             size=13, space_after=0)


# ---------------------------------------------------------------------------
@slide
def s_architecture(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Architecture", eyebrow="System overview")

    # Layered boxes
    layers = [
        ("CALLER", "society.run(handoff)", ACCENT_DARK, INK),
        ("HANDOFF ENVELOPE", "task + intent + assumptions + prior + confidence_required + human_review_when",
         ACCENT, RGBColor(0xFF, 0xFF, 0xFF)),
        ("ROUTER", "TF-IDF (or pluggable embed_fn) over skill index  →  ranked (agent, skill) candidates",
         ACCENT, RGBColor(0xFF, 0xFF, 0xFF)),
        ("CONFORMANCE GUARDRAIL", "declared-skill check  +  boundary allow/deny (NFKC-normalised)",
         GOOD, RGBColor(0xFF, 0xFF, 0xFF)),
        ("DISPATCHER", "LocalTransport (adapters)  ·  HttpTransport (A2A JSON-RPC)  ·  CompositeTransport",
         ACCENT_DARK, RGBColor(0xFF, 0xFF, 0xFF)),
        ("LEDGER + EXPLANATION STORE", "RoutingRecord + RoutingExplanation, keyed by handoff id  ·  JSONL or in-memory",
         MUTED, RGBColor(0xFF, 0xFF, 0xFF)),
        ("GOVERNANCE DETECTORS", "ConflictDetector  ·  CapabilityDriftDetector  ·  LowConfidence  ·  HumanReview  →  hooks (side-effect only)",
         WARN, RGBColor(0xFF, 0xFF, 0xFF)),
        ("OBSERVABILITY", "MetricsCollector (Prometheus)  ·  SessionTracer  ·  agent2society logger",
         INK, RGBColor(0xFF, 0xFF, 0xFF)),
    ]
    y = Inches(1.7)
    h = Inches(0.62)
    gap = Inches(0.08)
    for label, body, fill, fg in layers:
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, y, SLIDE_W - 2 * MARGIN, h)
        rect.adjustments[0] = 0.2
        rect.fill.solid(); rect.fill.fore_color.rgb = fill; rect.line.fill.background()
        rect.shadow.inherit = False
        # label
        box, tf = add_textbox(s, MARGIN + Inches(0.2), y + Inches(0.08), Inches(3.5), Inches(0.5))
        p = tf.paragraphs[0]
        r = p.add_run()
        set_run(r, text=label, font=HEAD, size=14, bold=True, color=fg)
        # body
        box, tf = add_textbox(s, MARGIN + Inches(3.8), y + Inches(0.1), SLIDE_W - 2 * MARGIN - Inches(4.0), Inches(0.5))
        p = tf.paragraphs[0]
        r = p.add_run()
        set_run(r, text=body, font=MONO, size=12, color=fg)
        y = y + h + gap


# ---------------------------------------------------------------------------
@slide
def s_api(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Sixty-second tour", eyebrow="API surface")

    code = (
        "from agent2society import Society, Handoff\n"
        "\n"
        "society = Society()\n"
        "society.add(\"https://research-agent.acme.com/.well-known/agent-card.json\")\n"
        "society.add(\"https://writer-agent.acme.com/.well-known/agent-card.json\")\n"
        "society.add(my_crewai_crew)                              # adapters wrap native objects\n"
        "\n"
        "society.boundary(\"writer-agent\", deny=[\"financial-data\"])\n"
        "\n"
        "# A Handoff carries why the task exists, what we assume, and prior decisions.\n"
        "h = Handoff(\n"
        "    task=\"Draft an executive memo on Q3 customer churn\",\n"
        "    intent=\"prep the board pack\",\n"
        "    assumptions=[\"churn data is final\"],\n"
        "    confidence_required=0.5,\n"
        ")\n"
        "memo = society.run(h)\n"
        "\n"
        "print(society.explain(h.id).render())   # human-readable rationale\n"
        "society.report()                        # per-task ledger"
    )
    add_code_block(s, code, MARGIN, Inches(1.7), Inches(8.3), Inches(5.2), size=12)

    box, tf = add_textbox(s, Inches(9.1), Inches(1.7), Inches(3.75), Inches(5.2))
    add_para(tf, "Five primitives", font=HEAD, size=14, bold=True, color=ACCENT, first=True, space_after=10)
    add_para(tf, "Society — the mesh", bullet=True, size=14, bold=True, space_after=4)
    add_para(tf, "Handoff — task + context", bullet=True, size=14, bold=True, space_after=4)
    add_para(tf, "boundary() — policy", bullet=True, size=14, bold=True, space_after=4)
    add_para(tf, "explain() — rationale", bullet=True, size=14, bold=True, space_after=4)
    add_para(tf, "report() — audit log", bullet=True, size=14, bold=True, space_after=14)
    add_para(tf, "Backward compat", font=HEAD, size=13, bold=True, color=MUTED, space_after=4)
    add_para(tf, "Mesh kept as alias for Society. Bare strings still work via Handoff.from_string().",
             size=12, color=MUTED, space_after=0)


# ---------------------------------------------------------------------------
@slide
def s_routing(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "How routing works", eyebrow="Routing")

    box, tf = add_textbox(s, MARGIN, Inches(1.7), Inches(6.0), Inches(5.3))
    add_para(tf, "Per (agent, skill) pair, the router scores:", font=HEAD, size=15, bold=True, first=True, space_after=10)
    add_para(tf, "Cosine similarity of skill text vs task text", bullet=True, size=15, space_after=6)
    add_para(tf, "Deterministic tag-overlap bonus", bullet=True, size=15, space_after=12)
    add_para(tf, "Returns a ranked list. Conformance runs on the top candidate; on failure it falls through to the next.",
             size=14, color=MUTED, space_after=10)
    add_para(tf, "The runner-up gets a rejected_reason recorded on the explanation — you see WHY the loser lost.",
             size=14, color=MUTED, italic=True, space_after=14)
    add_para(tf, "Swap the scorer", font=HEAD, size=14, bold=True, color=ACCENT, space_after=8)
    add_para(tf, "Default is dependency-free TF-IDF so transparency is visible on first run. embed_fn= takes any callable.",
             size=13, space_after=0)

    code = (
        "from sentence_transformers import SentenceTransformer\n"
        "from agent2society import Society\n"
        "\n"
        "model = SentenceTransformer(\"all-MiniLM-L6-v2\")\n"
        "society = Society(\n"
        "    embed_fn=lambda texts: model.encode(texts).tolist()\n"
        ")\n"
        "\n"
        "# Or supply OpenAI / Voyage / your hosted embedder:\n"
        "society = Society(embed_fn=my_remote_embedder)"
    )
    add_code_block(s, code, Inches(7.1), Inches(1.7), Inches(5.85), Inches(3.4), size=13)

    # Inline flags panel
    add_panel(s, Inches(7.1), Inches(5.25), Inches(5.85), Inches(1.7))
    box, tf = add_textbox(s, Inches(7.3), Inches(5.4), Inches(5.5), Inches(1.5))
    add_para(tf, "Quality flags on every explanation", font=HEAD, size=13, bold=True, color=ACCENT, first=True, space_after=6)
    add_para(tf, "OOD — top-1 far from any skill", bullet=True, size=12, space_after=3)
    add_para(tf, "VECTOR_AMBIGUITY — multiple near-ties", bullet=True, size=12, space_after=3)
    add_para(tf, "LOW_MARGIN — small top-1 vs top-2 gap", bullet=True, size=12, space_after=0)


# ---------------------------------------------------------------------------
@slide
def s_conformance(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Conformance: the guardrail", eyebrow="Safety")

    box, tf = add_textbox(s, MARGIN, Inches(1.7), Inches(6.0), Inches(5.3))
    add_para(tf, "Before dispatch, deterministically answer:", font=HEAD, size=16, bold=True, first=True, space_after=10)
    add_para(tf, "Does this agent's card actually declare this skill?", bullet=True, size=15, space_after=6)
    add_para(tf, "Is the task inside this agent's declared boundary (allow / deny)?", bullet=True, size=15, space_after=14)
    add_para(tf, "If either fails:", font=HEAD, size=14, bold=True, color=ACCENT, space_after=6)
    add_para(tf, "strict=True (default) — raises ConformanceViolation", bullet=True, size=14, space_after=6)
    add_para(tf, "strict=False — records in telemetry + explanation.blocked_reason, returns empty string",
             bullet=True, size=14, space_after=14)
    add_para(tf, "v0.5.3 hardening", font=HEAD, size=14, bold=True, color=GOOD, space_after=6)
    add_para(tf, "Boundary terms compared after NFKC + casefold. \"FINANCIAL\", \"fınancial\", \"ﬁnancial\" all collapse to the same key. Unicode bypass closed.",
             size=13, italic=True, color=MUTED, space_after=0)

    code = (
        "society.boundary(\n"
        "    \"writer-agent\",\n"
        "    deny=[\"financial-data\", \"pii\"]\n"
        ")\n"
        "society.boundary(\n"
        "    \"research-agent\",\n"
        "    allow=[\"public\"]\n"
        ")\n"
        "\n"
        "# Boundary edits are copy-on-write:\n"
        "# the graph is deep-copied, mutated,\n"
        "# and atomically swapped. Concurrent\n"
        "# run() never sees a partial state."
    )
    add_code_block(s, code, Inches(7.1), Inches(1.7), Inches(5.85), Inches(5.2), size=13)


# ---------------------------------------------------------------------------
@slide
def s_handoff(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Handoff: context that survives the chain", eyebrow="Meaning")

    box, tf = add_textbox(s, MARGIN, Inches(1.65), SLIDE_W - 2 * MARGIN, Inches(0.9))
    add_para(tf,
             "A bare string answers \"what to do next.\" A Handoff also answers why it exists, "
             "what we assume, what's been done so far, and what would force a human review.",
             font=HEAD, size=16, italic=True, color=INK, first=True, space_after=0)

    code = (
        "h0 = Handoff(\n"
        "    task=\"Research Q3 customer churn drivers\",\n"
        "    intent=\"prep the Q3 board pack\",\n"
        "    assumptions=[\"churn data through end-of-quarter is final\"],\n"
        "    confidence_required=0.5,\n"
        "    human_review_when=lambda r: \"cite\" not in r,\n"
        ")\n"
        "research = society.run(h0)\n"
        "\n"
        "# Extend the handoff — the memo dispatch sees the upstream decision\n"
        "# in its prior chain (and so does the routing explanation).\n"
        "h1 = h0.extend(\n"
        "    agent=\"research-agent\",\n"
        "    skill=\"web_research\",\n"
        "    summary=\"found 3 churn drivers\",\n"
        "    confidence=0.62,\n"
        "    next_task=\"Draft an executive memo on those churn drivers\",\n"
        ")\n"
        "memo = society.run(h1)"
    )
    add_code_block(s, code, MARGIN, Inches(2.8), Inches(8.3), Inches(4.1), size=13)

    box, tf = add_textbox(s, Inches(9.1), Inches(2.8), Inches(3.75), Inches(4.1))
    add_para(tf, "Why it matters", font=HEAD, size=14, bold=True, color=ACCENT, first=True, space_after=8)
    add_para(tf, "Ledger becomes a directed graph of decisions, not independent rows.",
             bullet=True, size=13, space_after=8)
    add_para(tf, "Downstream explanations carry the upstream confidence and rationale.",
             bullet=True, size=13, space_after=8)
    add_para(tf, "human_review_when default is fail-safe — exceptions trigger review.",
             bullet=True, size=13, space_after=0)


# ---------------------------------------------------------------------------
@slide
def s_selfassessment(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "SelfAssessment: agents declare their own limits", eyebrow="Meaning")

    code = (
        '{\n'
        '  "name": "writer-agent",\n'
        '  "skills": [\n'
        '    {"id": "exec_memo", "name": "Executive Memo"}\n'
        '  ],\n'
        '  "selfAssessment": {\n'
        '    "confidenceModel": "tfidf_score",\n'
        '    "knownLimitations": [\n'
        '      "English only",\n'
        '      "max ~400 words per memo"\n'
        '    ],\n'
        '    "outOfScope": [\n'
        '      "legal opinion",\n'
        '      "binding financial guidance"\n'
        '    ],\n'
        '    "escalateWhen": [\n'
        '      "any quantitative claim cited without source"\n'
        '    ]\n'
        '  }\n'
        '}'
    )
    add_code_block(s, code, MARGIN, Inches(1.7), Inches(7.0), Inches(5.2), size=13)

    box, tf = add_textbox(s, Inches(7.85), Inches(1.7), Inches(5.0), Inches(5.2))
    add_para(tf, "First-class caveats", font=HEAD, size=14, bold=True, color=ACCENT, first=True, space_after=8)
    add_para(tf,
             "An agent's own scope statement is the most reliable place to put it.",
             size=14, space_after=10)
    add_para(tf,
             "agent2society surfaces these limits on every explanation that picks the agent — caller AND downstream agents see the same caveats the agent claims for itself.",
             size=13, color=MUTED, space_after=14)
    add_para(tf, "Supervisor pattern loses this", font=HEAD, size=14, bold=True, color=BAD, space_after=6)
    add_para(tf,
             "Caveats live inside an agent's prompt where the caller never sees them. agent2society treats them as decision metadata.",
             size=13, color=MUTED, space_after=0)


# ---------------------------------------------------------------------------
@slide
def s_governance(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Governance hooks: detection, never auto-correct", eyebrow="Governance")

    code = (
        "society.on_low_confidence(\n"
        "    lambda exp: notify(exp), threshold=0.5\n"
        ")\n"
        "society.on_human_review(\n"
        "    lambda exp, result: page_oncall(result)\n"
        ")\n"
        "society.on_conflict(\n"
        "    lambda c: log(\"conflict\", c.detail)\n"
        ")\n"
        "society.on_capability_drift(\n"
        "    lambda d: log(\"drift\", d.agent)\n"
        ")"
    )
    add_code_block(s, code, MARGIN, Inches(1.7), Inches(6.0), Inches(3.6), size=13)

    box, tf = add_textbox(s, Inches(6.85), Inches(1.7), Inches(6.0), Inches(3.6))
    add_para(tf, "on_low_confidence", font=HEAD, size=13, bold=True, color=ACCENT, first=True, space_after=2)
    add_para(tf, "decision below handoff's confidence_required (or society threshold)",
             size=12, color=MUTED, space_after=8)
    add_para(tf, "on_human_review", font=HEAD, size=13, bold=True, color=ACCENT, space_after=2)
    add_para(tf, "Handoff.human_review_when(result) returned True", size=12, color=MUTED, space_after=8)
    add_para(tf, "on_conflict", font=HEAD, size=13, bold=True, color=ACCENT, space_after=2)
    add_para(tf, "same task text routed to different (agent, skill) pairs", size=12, color=MUTED, space_after=8)
    add_para(tf, "on_capability_drift", font=HEAD, size=13, bold=True, color=ACCENT, space_after=2)
    add_para(tf, "one agent winning across an unusually broad spread of skills",
             size=12, color=MUTED, space_after=0)

    # Bottom callout
    add_panel(s, MARGIN, Inches(5.55), SLIDE_W - 2 * MARGIN, Inches(1.4), fill=RGBColor(0xFF, 0xF7, 0xE6))
    box, tf = add_textbox(s, MARGIN + Inches(0.2), Inches(5.7), SLIDE_W - 2 * MARGIN - Inches(0.4), Inches(1.1))
    add_para(tf, "Hooks are side effects. They cannot block, retry, or modify a dispatch.",
             font=HEAD, size=15, bold=True, color=WARN, first=True, space_after=6)
    add_para(tf,
             "Handler exceptions are caught (a buggy hook can never break a dispatch) but logged at WARNING with the hook's qualname. "
             "Silent auto-correction is exactly the opacity that makes meshes untrustworthy.",
             size=13, color=INK, space_after=0)


# ---------------------------------------------------------------------------
@slide
def s_quality_signals(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Routing-quality signals you should actually watch", eyebrow="Observability")

    cards = [
        ("OOD",
         "Out-of-distribution.",
         "Top-1 task is far from anything in the skill index. Best guess returned, but should be treated as low quality.",
         BAD),
        ("VECTOR_AMBIGUITY",
         "Multiple near-ties.",
         "Skill descriptions overlap too much, or this task genuinely has two valid handlers.",
         WARN),
        ("LOW_MARGIN",
         "Small top-1 vs top-2 gap.",
         "Cheapest confidence proxy. Surface on dashboards, alert on dips, or gate retry behavior on it.",
         ACCENT),
    ]
    card_w = (SLIDE_W - 2 * MARGIN - Inches(0.4)) / 3
    card_h = Inches(3.2)
    y = Inches(1.85)
    x = MARGIN
    for tag, oneliner, body, color in cards:
        add_panel(s, x, y, card_w, card_h)
        # Top stripe + tag
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.18), y + Inches(0.15), card_w - Inches(0.36), Inches(0.08))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = color; stripe.line.fill.background()
        box, tf = add_textbox(s, x + Inches(0.25), y + Inches(0.35), card_w - Inches(0.5), card_h - Inches(0.6))
        add_para(tf, tag, font=MONO, size=18, bold=True, color=color, first=True, space_after=6)
        add_para(tf, oneliner, font=HEAD, size=14, bold=True, color=INK, space_after=10)
        add_para(tf, body, size=13, color=MUTED, space_after=0)
        x += card_w + Inches(0.2)

    add_panel(s, MARGIN, Inches(5.3), SLIDE_W - 2 * MARGIN, Inches(1.65), fill=BG_PANEL)
    box, tf = add_textbox(s, MARGIN + Inches(0.2), Inches(5.45), SLIDE_W - 2 * MARGIN - Inches(0.4), Inches(1.4))
    add_para(tf, "Feedback loop", font=HEAD, size=14, bold=True, color=ACCENT, first=True, space_after=4)
    add_para(tf,
             "Low-margin decisions feed a human-review queue. The queue's outputs feed your skill descriptions. "
             "Your skill descriptions feed back into the router. The mesh improves without a new model in the critical path.",
             size=13, color=INK, space_after=0)


# ---------------------------------------------------------------------------
@slide
def s_prod_ops(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Production ops", eyebrow="Run it in real systems")

    items = [
        ("Metrics",
         "Prometheus text or JSON snapshot.",
         "society.metrics.render_prometheus()",
         ACCENT),
        ("Persistent ledger",
         "JSONL store — append-on-run, rebuild-on-restart, corruption-tracked.",
         "Society(store=JsonlFileStore(\"audit.jsonl\"))",
         ACCENT_DARK),
        ("Auto-retry / fallback",
         "Transport errors fall through to next conformance-passing candidate.",
         "society.run(handoff, retry=True)",
         GOOD),
        ("Thread safety",
         "Reentrant lock + copy-on-write boundary edits. Concurrent run() is safe.",
         "society.add(...) / society.run(...) / society.optimize(...)",
         WARN),
        ("Session traces",
         "Joins telemetry + explanation store into an ordered event stream.",
         "for ev in SessionTracer(society).events(): ...",
         ACCENT),
        ("LLM-assisted optimizer",
         "LLM proposes tokens; backtest decides. LLM is NEVER in the routing critical path.",
         "society.optimize(labels, llm_fn=my_proposer)",
         ACCENT_DARK),
    ]
    cols = 2
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / cols
    row_h = Inches(1.55)
    for i, (title, body, code, color) in enumerate(items):
        col = i % cols
        row = i // cols
        x = MARGIN + col * (col_w + Inches(0.3))
        y = Inches(1.75) + row * (row_h + Inches(0.15))
        add_panel(s, x, y, col_w, row_h)
        # color bar
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.12), Inches(0.08), row_h - Inches(0.24))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = color; stripe.line.fill.background()
        box, tf = add_textbox(s, x + Inches(0.25), y + Inches(0.1), col_w - Inches(0.4), row_h - Inches(0.2))
        add_para(tf, title, font=HEAD, size=14, bold=True, color=color, first=True, space_after=4)
        add_para(tf, body, size=12, color=INK, space_after=6)
        add_para(tf, code, font=MONO, size=11, color=MUTED, space_after=0)


# ---------------------------------------------------------------------------
@slide
def s_adapters(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Native agents: adapters", eyebrow="Framework integration")

    box, tf = add_textbox(s, MARGIN, Inches(1.7), SLIDE_W - 2 * MARGIN, Inches(0.7))
    add_para(tf,
             "Adapters wrap native objects into A2A cards + local handlers. "
             "CompositeTransport prefers local handlers over HTTP — no network in the loop.",
             font=HEAD, size=15, italic=True, color=INK, first=True, space_after=0)

    # Logos / framework chips
    chips = [
        ("CrewAI", ".kickoff() / .agents", ACCENT),
        ("LangGraph", ".invoke() / .nodes", ACCENT_DARK),
        ("AutoGen", "ConversableAgent / GroupChatManager", GOOD),
        ("Callables", "any def / lambda", MUTED),
        ("Duck-typed", "run / invoke / kickoff", WARN),
    ]
    x = MARGIN
    y = Inches(2.7)
    for name, sub, color in chips:
        w = Inches(2.4)
        h = Inches(1.4)
        add_panel(s, x, y, w, h)
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.18), y + Inches(0.18), w - Inches(0.36), Inches(0.08))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = color; stripe.line.fill.background()
        box, tf = add_textbox(s, x + Inches(0.25), y + Inches(0.38), w - Inches(0.5), h - Inches(0.5))
        add_para(tf, name, font=HEAD, size=16, bold=True, color=INK, first=True, space_after=4)
        add_para(tf, sub, font=MONO, size=11, color=MUTED, space_after=0)
        x += w + Inches(0.1)

    code = (
        "from agent2society.adapters.base import Adapter, register_adapter\n"
        "\n"
        "class MyAdapter(Adapter):\n"
        "    def matches(self, obj): ...     # claim this object?\n"
        "    def to_card(self, obj): ...     # build an A2A card\n"
        "    def to_handler(self, obj): ...  # local-dispatch callable\n"
        "\n"
        "register_adapter(MyAdapter())\n"
        "society.add(my_native_thing)        # routes the same as any HTTP agent"
    )
    add_code_block(s, code, MARGIN, Inches(4.4), SLIDE_W - 2 * MARGIN, Inches(2.45), size=13)


# ---------------------------------------------------------------------------
@slide
def s_v053(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "v0.5.3: enterprise hardening", eyebrow="What's new")

    box, tf = add_textbox(s, MARGIN, Inches(1.65), SLIDE_W - 2 * MARGIN, Inches(0.55))
    add_para(tf,
             "15 production failure modes audited and fixed. 112 tests passing — 16 dedicated to regression on these specific failures.",
             font=HEAD, size=14, italic=True, color=MUTED, first=True, space_after=0)

    items = [
        ("Copy-on-write graph mutation",
         "Boundary edits deep-copy the graph, mutate the copy, atomic swap. Concurrent run() never sees a partial state."),
        ("Unicode bypass closed",
         "Conformance compares with NFKC + casefold. \"FINANCIAL\", \"fınancial\", \"ﬁnancial\" all collapse to the same key."),
        ("Governance memory bounded",
         "Conflict/drift detectors use FIFO-evicted dicts (10k cap). No unbounded leaks under sustained traffic."),
        ("Fail-safe defaults",
         "human_review_when exceptions default to True. Predicate failures get MORE reviews, not silently fewer."),
        ("Hook exceptions surfaced",
         "Caught (won't crash dispatch) but logged at WARNING with the hook's qualname and exception class."),
        ("Total response extraction",
         "extract_text always returns a string. Falls back to json.dumps(default=str), then repr()."),
        ("JSONL corruption stats",
         "skipped_json / skipped_shape tracked and logged. No more silent data loss on store rebuild."),
        ("Double-checked router rebuild",
         "Race-free lazy index build under concurrent first-access. No duplicate work, no half-built index visible."),
        ("Retry attribution",
         "Every failed attempt recorded in RoutingRecord.fallbacks + dispatch_retries_total counter."),
        ("DispatchError wrapping",
         "Handler exceptions wrapped once, no double-wrap. Preserves original cause."),
        ("Frozen snapshot mappings",
         "as_mapping() returns a point-in-time deep copy. Callers can't mutate the live store."),
        ("Adapter skips logged",
         "DEBUG-level log on every adapter that declines to match — diagnostics, not silence."),
    ]
    cols = 2
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / cols
    row_h = Inches(0.78)
    y_start = Inches(2.35)
    for i, (title, body) in enumerate(items):
        col = i % cols
        row = i // cols
        x = MARGIN + col * (col_w + Inches(0.3))
        y = y_start + row * (row_h + Inches(0.07))
        box, tf = add_textbox(s, x, y, col_w, row_h)
        add_para(tf, title, font=HEAD, size=12, bold=True, color=ACCENT, first=True, space_after=2)
        add_para(tf, body, size=11, color=INK, space_after=0)


# ---------------------------------------------------------------------------
@slide
def s_when_to(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "When to use, when not to", eyebrow="Honest positioning")

    # Use it when ...
    add_panel(s, MARGIN, Inches(1.75), Inches(6.05), Inches(5.1), fill=RGBColor(0xEC, 0xF7, 0xF0))
    box, tf = add_textbox(s, MARGIN + Inches(0.25), Inches(1.9), Inches(5.75), Inches(4.9))
    add_para(tf, "USE agent2society WHEN", font=HEAD, size=14, bold=True, color=GOOD, first=True, space_after=10)
    add_para(tf, "Multiple specialist agents with overlapping capabilities and you need defensible routing.",
             bullet=True, size=14, space_after=8)
    add_para(tf, "Audit requirements: \"why was this task sent here?\" must have an answer.",
             bullet=True, size=14, space_after=8)
    add_para(tf, "Cost-sensitive coordination at scale (kill the supervisor LLM round).",
             bullet=True, size=14, space_after=8)
    add_para(tf, "Boundary policy matters (PII, financial, regulatory).",
             bullet=True, size=14, space_after=8)
    add_para(tf, "You want detection of conflict / drift / low confidence as first-class signals.",
             bullet=True, size=14, space_after=0)

    # Don't use when ...
    add_panel(s, Inches(6.95), Inches(1.75), Inches(5.9), Inches(5.1), fill=RGBColor(0xFB, 0xF1, 0xF1))
    box, tf = add_textbox(s, Inches(7.15), Inches(1.9), Inches(5.6), Inches(4.9))
    add_para(tf, "DON'T REACH FOR IT WHEN", font=HEAD, size=14, bold=True, color=BAD, first=True, space_after=10)
    add_para(tf, "Single-agent tool-use loop. Pick a framework, skip this layer.",
             bullet=True, size=14, space_after=8)
    add_para(tf, "Open-ended planning / decomposition is the bottleneck. v1 ships no planner by design.",
             bullet=True, size=14, space_after=8)
    add_para(tf, "Agents are dynamically created per-task (no stable cards to embed).",
             bullet=True, size=14, space_after=8)
    add_para(tf, "Routing genuinely needs negotiation between agents (not just a match).",
             bullet=True, size=14, space_after=8)
    add_para(tf, "You actively want a supervisor that does reasoning between turns.",
             bullet=True, size=14, space_after=0)


# ---------------------------------------------------------------------------
@slide
def s_install(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Install and quick start", eyebrow="Get going")

    code = (
        "# Core, zero hard deps\n"
        "pip install agent2society\n"
        "\n"
        "# HTTP dispatch (httpx)\n"
        "pip install \"agent2society[http]\"\n"
        "\n"
        "# Stronger routing (sentence-transformers)\n"
        "pip install \"agent2society[embeddings]\"\n"
        "\n"
        "# Benchmarks (tiktoken + langgraph + langchain-core)\n"
        "pip install \"agent2society[bench]\""
    )
    add_code_block(s, code, MARGIN, Inches(1.75), Inches(6.2), Inches(3.7), size=14)

    box, tf = add_textbox(s, Inches(7.05), Inches(1.75), Inches(5.8), Inches(5.1))
    add_para(tf, "Resources", font=HEAD, size=15, bold=True, color=ACCENT, first=True, space_after=10)
    add_para(tf, "PyPI", font=HEAD, size=13, bold=True, space_after=2)
    add_para(tf, "pypi.org/project/agent2society", font=MONO, size=13, color=ACCENT, space_after=10)
    add_para(tf, "Source", font=HEAD, size=13, bold=True, space_after=2)
    add_para(tf, "github.com/graxella/agent2society", font=MONO, size=13, color=ACCENT, space_after=10)
    add_para(tf, "Examples", font=HEAD, size=13, bold=True, space_after=2)
    add_para(tf, "examples/basic_mesh.py", font=MONO, size=12, color=MUTED, space_after=2)
    add_para(tf, "examples/transparent_mesh.py", font=MONO, size=12, color=MUTED, space_after=10)
    add_para(tf, "Benchmarks", font=HEAD, size=13, bold=True, space_after=2)
    add_para(tf, "benchmarks/run.py            (synthetic)", font=MONO, size=12, color=MUTED, space_after=2)
    add_para(tf, "benchmarks/run_langgraph.py  (real LangGraph)", font=MONO, size=12, color=MUTED, space_after=0)

    # First-run nudge across the bottom
    add_panel(s, MARGIN, Inches(5.75), SLIDE_W - 2 * MARGIN, Inches(1.1), fill=RGBColor(0xEC, 0xF1, 0xFF))
    box, tf = add_textbox(s, MARGIN + Inches(0.25), Inches(5.9), SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(0.9))
    add_para(tf, "First run", font=HEAD, size=13, bold=True, color=ACCENT, first=True, space_after=4)
    add_para(tf,
             "Default TF-IDF scorer ships with the package — no model download, no network. "
             "Run examples/basic_mesh.py to see a real RoutingExplanation in your terminal in under a minute.",
             size=13, color=INK, space_after=0)


# ---------------------------------------------------------------------------
@slide
def s_close(prs):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, BG)

    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SLIDE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

    box, tf = add_textbox(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(1.5))
    add_para(tf, "Transparency is a property of the system,",
             font=HEAD, size=36, bold=True, color=INK, first=True, space_after=4)
    add_para(tf, "not a promise from an LLM.",
             font=HEAD, size=36, bold=True, color=ACCENT, space_after=0)

    box, tf = add_textbox(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(1.6))
    add_para(tf,
             "agent2society replaces opaque supervisor reasoning with a deterministic graph, "
             "a cheap embedding lookup, and a first-class explanation of every decision.",
             font=HEAD, size=18, italic=True, color=MUTED, first=True, space_after=8)
    add_para(tf,
             "When the routing genuinely needs richer semantics, swap the scorer. Keep every other property.",
             font=HEAD, size=16, color=INK, space_after=0)

    # Final code line
    add_code_block(s, "pip install agent2society", Inches(0.9), Inches(5.5), Inches(6.0), Inches(0.7), size=18)

    box, tf = add_textbox(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6))
    add_para(tf, "Apache-2.0   ·   github.com/graxella/agent2society   ·   pypi.org/project/agent2society",
             font=HEAD, size=12, color=MUTED, first=True, space_after=0)


# ============================================================================
# BUILD
# ============================================================================
for builder in _slide_builders:
    builder(prs)

total = len(prs.slides)
for i, sld in enumerate(prs.slides):
    if i == 0 or i == total - 1:
        # title and closing slides skip the footer/number bar
        continue
    add_footer(sld, i + 1, total)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"wrote {OUT}  ({total} slides)")
